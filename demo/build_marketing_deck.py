"""Generate the marketing / demo deck for the Goodcare Dental booking system.

Audience: clinic owners and staff — not engineers. Zero jargon. Plain
Malaysian English. Every feature is framed as "no more X, now you Y".

Run:  python3 demo/build_marketing_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Brand
BRAND = RGBColor(0x0D, 0x94, 0x88)
BRAND_DARK = RGBColor(0x0F, 0x76, 0x6E)
BRAND_LIGHT = RGBColor(0xCC, 0xFB, 0xF1)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x44, 0x40, 0x3D)
MUTED = RGBColor(0x78, 0x71, 0x6C)
PAPER = RGBColor(0xFA, 0xFA, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_BG = RGBColor(0xD1, 0xFA, 0xE5)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
RED = RGBColor(0xC0, 0x39, 0x2B)
RED_BG = RGBColor(0xFE, 0xE2, 0xE2)


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Helvetica"
    return box


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=INK, bullet_color=BRAND, gap=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        b = p.add_run()
        b.text = "● "
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        t = p.add_run()
        t.text = item
        t.font.size = Pt(size)
        t.font.color.rgb = color
        t.font.name = "Helvetica"
        p.space_after = Pt(gap)


def add_brand_bar(slide, sw):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()


def add_footer(slide, sw, sh, n, total):
    add_text(slide, Inches(0.5), sh - Inches(0.4), Inches(6), Inches(0.3),
             "Goodcare Dental — Clinic Booking System", size=9, color=MUTED)
    add_text(slide, sw - Inches(1.5), sh - Inches(0.4), Inches(1), Inches(0.3),
             f"{n} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_section_label(slide, left, top, text, color=BRAND):
    """Small uppercase section pill above the slide title."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, Inches(2.2), Inches(0.36))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    tf = rect.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Helvetica"


def add_card(slide, left, top, width, height, title, body, *,
             title_color=BRAND_DARK, fill=WHITE, border=BRAND):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = border
    rect.line.width = Pt(1.5)

    add_text(slide, left + Inches(0.25), top + Inches(0.2),
             width - Inches(0.5), Inches(0.5),
             title, size=15, bold=True, color=title_color)
    add_text(slide, left + Inches(0.25), top + Inches(0.75),
             width - Inches(0.5), height - Inches(1),
             body, size=11, color=INK_SOFT)


def add_screenshot_box(slide, left, top, width, height, label):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = PAPER
    rect.line.color.rgb = MUTED
    rect.line.width = Pt(0.75)
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"📷  {label}"
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED
    r.font.italic = True
    r.font.name = "Helvetica"


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# -----------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

TOTAL = 22


def new_slide():
    s = prs.slides.add_slide(BLANK)
    add_brand_bar(s, SW)
    return s


# =================================================================
# 1. TITLE
# =================================================================
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.18), SW, SH - Inches(0.18))
bg.fill.solid()
bg.fill.fore_color.rgb = PAPER
bg.line.fill.background()

add_text(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.6),
         "GOODCARE DENTAL", size=14, color=BRAND_DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2),
         "Run your clinic without the headache.",
         size=44, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.6),
         "Patients book online. Your nurse stays in control.",
         size=20, color=INK_SOFT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.6),
         "WhatsApp does the talking. No more phone tag, no more lost appointments.",
         size=20, color=INK_SOFT, align=PP_ALIGN.CENTER)

# Brand "leaf" graphic-ish placeholder
oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.85), Inches(5.4),
                          Inches(1.6), Inches(1.6))
oval.fill.solid()
oval.fill.fore_color.rgb = BRAND
oval.line.fill.background()
add_text(s, Inches(5.85), Inches(5.9), Inches(1.6), Inches(0.6),
         "🦷", size=42, align=PP_ALIGN.CENTER)

add_notes(s,
    "Open warmly. 30 seconds tops.\n\n"
    "'Today I'll walk you through a system I built specifically for clinics like yours. "
    "Nothing fancy — just the boring problems solved. Phone bookings, reminders, no-shows, "
    "doctor schedules. By the end you'll see exactly how it would fit into your week.'\n\n"
    "Don't promise anything. Just set the frame: 'It's built for clinics, not for software people.'")


# =================================================================
# 2. THE PROBLEM
# =================================================================
s = new_slide()
add_text(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6),
         "Today's clinic, today's pain", size=28, bold=True, color=INK)
add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
         "If any of these sound familiar, this system is for you.",
         size=14, color=MUTED)

pains = [
    ("📞", "Phone rings during treatment",
     "Nurse picks up, doctor waits. Patient on the line asks for next Tuesday — but the diary is on the other table."),
    ("📓", "Paper diary or shared notebook",
     "Two patients booked for 3pm. Nurse A wrote one, Nurse B wrote another. Now you have an awkward apology."),
    ("⏰", "Patients forget. They don't come.",
     "No-show after no-show. You lose RM 100, RM 150 each time. Nobody remembered to remind them."),
    ("📋", "Who confirmed Mrs Lim's booking?",
     "Patient calls asking about her appointment. Three nurses, four notebooks. Nobody knows for sure."),
    ("🕘", "Doctor on leave next week",
     "Receptionist forgot. Patient booked. Now you call them back to apologise."),
    ("📊", "How busy was last month?",
     "You want to know — but the data lives in WhatsApp chats and scribbles. You'd need a weekend to count."),
]
gap = Inches(0.2)
card_w = Inches(4.0)
card_h = Inches(1.7)
xs = [Inches(0.8), Inches(0.8) + card_w + gap, Inches(0.8) + (card_w + gap) * 2]
ys = [Inches(2.0), Inches(2.0) + card_h + gap]
for i, (icon, title, body) in enumerate(pains):
    row = i // 3
    col = i % 3
    x = xs[col]
    y = ys[row]
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = RED
    rect.line.width = Pt(1)
    add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(0.6), Inches(0.5),
             icon, size=22)
    add_text(s, x + Inches(0.95), y + Inches(0.2), card_w - Inches(1.2), Inches(0.5),
             title, size=13, bold=True, color=INK)
    add_text(s, x + Inches(0.25), y + Inches(0.8), card_w - Inches(0.5),
             card_h - Inches(0.9), body, size=10, color=INK_SOFT)

add_footer(s, SW, SH, 2, TOTAL)
add_notes(s,
    "Pause on each card. Let them recognise the pain. Don't ask if it happens — they know it does.\n\n"
    "Your tone: 'You've lived this. I built this because clinics shouldn't have to.'\n\n"
    "If you only have a minute on this slide: just read the headlines. 'Phone rings during treatment. "
    "Two patients booked for 3pm. Doctor on leave but receptionist forgot.' Then move on.")


# =================================================================
# 3. THE PROMISE
# =================================================================
s = new_slide()
add_text(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6),
         "Here's what changes the first week", size=28, bold=True, color=INK)
add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
         "Three promises. Each one is something you can verify on day one.",
         size=14, color=MUTED)

promises = [
    ("Patients book themselves",
     "They tap a link, fill a form, pick a slot. Your nurse never picks up an appointment call again.",
     "🌐"),
    ("Your nurse stays in charge",
     "Every booking still passes through her. Approve, reschedule, cancel — one tap each. WhatsApp messages pre-typed.",
     "👩‍⚕️"),
    ("Nothing falls through the cracks",
     "Reminders the day before. Track who came and who didn't. See every action with a name and time.",
     "✅"),
]
y = Inches(2.5)
card_w = Inches(3.8)
gap = Inches(0.3)
total_w = card_w * 3 + gap * 2
start_x = (SW - total_w) / 2
for i, (title, body, icon) in enumerate(promises):
    x = start_x + (card_w + gap) * i
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(3.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BRAND
    rect.line.width = Pt(2)
    # icon circle
    cir = s.shapes.add_shape(MSO_SHAPE.OVAL, x + card_w/2 - Inches(0.6),
                             y + Inches(0.4), Inches(1.2), Inches(1.2))
    cir.fill.solid()
    cir.fill.fore_color.rgb = BRAND_LIGHT
    cir.line.fill.background()
    add_text(s, x + card_w/2 - Inches(0.6), y + Inches(0.6),
             Inches(1.2), Inches(0.8), icon, size=36, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(1.85),
             card_w - Inches(0.5), Inches(0.5),
             title, size=16, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(2.4),
             card_w - Inches(0.5), Inches(1.0),
             body, size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)

add_footer(s, SW, SH, 3, TOTAL)
add_notes(s,
    "These are the only three things you need them to remember.\n\n"
    "Land each one with conviction:\n"
    "1. 'Your nurse doesn't take booking calls anymore.'\n"
    "2. 'Every booking still passes through her — she's still the gatekeeper.'\n"
    "3. 'You'll know exactly who showed up, who reminded them, who approved what.'\n\n"
    "Don't sell features yet. Sell outcomes.")


# =================================================================
# 4. MEET THE TEAM
# =================================================================
s = new_slide()
add_text(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6),
         "Four people, four different views", size=28, bold=True, color=INK)
add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
         "Each person logs in and sees only what they need. No clutter.",
         size=14, color=MUTED)

people = [
    ("🧑‍🦰", "Patient", "The simplest. Just fills a form. No login, no app to install.",
     "Books, reschedules, or cancels online. Done in under a minute."),
    ("👩‍⚕️", "Nurse / Receptionist", "The daily driver. Most of her work happens here.",
     "Approves bookings, sends WhatsApp, marks attendance, handles walk-ins."),
    ("🦷", "Doctor", "Just sees their own day.", "Today's appointments, this week's calendar, time to block off (lunch, MC)."),
    ("👨‍💼", "Owner (you)", "Sees everything. Edits anything. Holds everyone accountable.",
     "Reviews bookings, approves leave, edits message templates, audits every action."),
]
y = Inches(2.4)
card_w = Inches(2.95)
gap = Inches(0.15)
total_w = card_w * 4 + gap * 3
start_x = (SW - total_w) / 2
for i, (icon, role, sub, body) in enumerate(people):
    x = start_x + (card_w + gap) * i
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(3.7))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = MUTED
    rect.line.width = Pt(0.75)
    add_text(s, x, y + Inches(0.3), card_w, Inches(0.8),
             icon, size=44, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.4),
             Inches(0.4), role, size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.8), card_w - Inches(0.4),
             Inches(0.5), sub, size=11, color=BRAND_DARK, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, x + Inches(0.2), y + Inches(2.4), card_w - Inches(0.4),
             Inches(1.2), body, size=11, color=INK_SOFT, align=PP_ALIGN.CENTER)

add_footer(s, SW, SH, 4, TOTAL)
add_notes(s,
    "Walk left-to-right. Each card is one sentence.\n\n"
    "If they ask 'what about the cleaner / accountant / etc.' — answer: 'We can add more roles "
    "later. For now: patient, nurse, doctor, owner. That covers 95% of clinic work.'")


# =================================================================
# Helper: section divider slide
# =================================================================
def section_divider(label, headline, subline, page_num):
    s = new_slide()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.18),
                            SW, SH - Inches(0.18))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND
    bg.line.fill.background()

    add_text(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5),
             label.upper(), size=14, bold=True, color=BRAND_LIGHT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.2),
             headline, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.6),
             subline, size=18, color=BRAND_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, SW, SH, page_num, TOTAL)
    return s


# Helper: standard feature slide
def feature_slide(section, title, sub, bullets, screenshot_label,
                  notes_text, page_num):
    s = new_slide()
    add_section_label(s, Inches(0.8), Inches(0.55), section)
    add_text(s, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.7),
             title, size=28, bold=True, color=INK)
    add_text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.5),
             sub, size=14, color=MUTED, italic=True)

    add_text(s, Inches(0.8), Inches(2.55), Inches(5.7), Inches(0.4),
             "What it does for you", size=12, bold=True, color=BRAND_DARK)
    add_bullets(s, Inches(0.8), Inches(2.95), Inches(5.7), Inches(3.5),
                bullets, size=13)

    add_text(s, Inches(6.85), Inches(2.55), Inches(5.7), Inches(0.4),
             "What it looks like", size=12, bold=True, color=BRAND_DARK)
    add_screenshot_box(s, Inches(6.85), Inches(2.95),
                       Inches(5.7), Inches(3.7), screenshot_label)

    add_footer(s, SW, SH, page_num, TOTAL)
    add_notes(s, notes_text)
    return s


# =================================================================
# 5. PATIENT DIVIDER
# =================================================================
section_divider("Part 1 of 4", "What your patient sees",
                "One link. One form. Sixty seconds.", 5)

# =================================================================
# 6. Patient booking form
# =================================================================
feature_slide(
    section="The patient",
    title="A single form. No app to download.",
    sub="Send them a link. That's it. Works on any phone, any computer.",
    bullets=[
        "Patient fills name, IC or passport, WhatsApp number",
        "Picks the treatment — scaling, root canal, whitening, wisdom tooth, or 'others'",
        "Sees only times when the doctor is actually free",
        "Picks one, hits Submit. Done.",
        "Gets a message: 'Our nurse will contact you within 24 hours.'",
    ],
    screenshot_label="The /book form on a phone screen",
    notes_text=(
        "Show this on the projector if you can. Open /book and let them see it.\n\n"
        "Key things to highlight while pointing:\n"
        "- 'See how the country sets the phone code automatically. Malaysian patient = +60 pre-filled. "
        "They just type the local digits.'\n"
        "- 'The treatment they pick controls how long the slot is. A scaling is 30 min, root canal 90. "
        "That way the calendar fills up cleanly.'\n"
        "- 'Date picker only allows future dates. Sundays show no slots if the doctor doesn't work "
        "Sundays — they can't even pick that.'"
    ),
    page_num=6,
)

# =================================================================
# 7. Smart slot picker
# =================================================================
feature_slide(
    section="The patient",
    title="The slot picker is smart",
    sub="No more 'sorry, that doctor is on leave' phone calls.",
    bullets=[
        "Doctor on leave that day? Those dates don't appear at all.",
        "Doctor only works half-day Wednesday? Only morning slots show.",
        "Lunch break? Hidden automatically.",
        "Someone else already booked that 10am slot? It disappears the second they confirm.",
        "Treatment takes 90 minutes? Only slots with 90 free minutes show up.",
    ],
    screenshot_label="Slot picker with greyed-out + available times",
    notes_text=(
        "This is where you separate yourself from Calendly et al.\n\n"
        "'Generic booking tools don't know your doctors are on leave on Tuesday. They show every "
        "Tuesday. Patient picks one. You scramble. We know — because leave is built into the system. "
        "Approve a leave request once, and the public form automatically blocks those days.'\n\n"
        "Mention the smart gap-filling: 'If two existing appointments leave a 60-minute gap in the "
        "middle of the day, a 60-min patient sees that gap. Nothing wasted.'"
    ),
    page_num=7,
)

# =================================================================
# 8. Reschedule + cancel
# =================================================================
feature_slide(
    section="The patient",
    title="Reschedule or cancel — same form, same link",
    sub="Patient doesn't need to call to change their appointment.",
    bullets=[
        "Same /book page. Tabs at the top: Booking · Reschedule · Cancel.",
        "Patient enters their IC and phone. System finds their existing booking.",
        "Reschedule: pick a new time. Cancel: confirm once.",
        "Both go to your nurse for approval — never automatic, you stay in control.",
        "Old slot gets freed up the moment the nurse approves the reschedule.",
    ],
    screenshot_label="Reschedule tab with patient's active booking shown",
    notes_text=(
        "'Patient changes their mind, family thing comes up. They don't have to call. They go to the "
        "same form, pick Reschedule, find their booking, pick a new time. Your nurse sees it as a "
        "request — same flow as a new booking.'\n\n"
        "Why this matters: 'It removes one of the most annoying calls your nurse gets. And the patient "
        "never feels guilty about asking — they just click.'"
    ),
    page_num=8,
)


# =================================================================
# 9. NURSE DIVIDER
# =================================================================
section_divider("Part 2 of 4", "What your nurse does each day",
                "Approve, message, remind, track. All in one screen.", 9)


# =================================================================
# 10. Pending queue
# =================================================================
feature_slide(
    section="The nurse",
    title="One screen shows everything waiting for her",
    sub="No more checking three different places. The 'Pending' queue is the morning routine.",
    bullets=[
        "Every new booking, reschedule request, or cancellation lands here",
        "Each card shows: name, IC, WhatsApp, doctor, time, treatment",
        "First-time patient gets a purple 'First-time' badge",
        "Returning patient shows '3 prior visits' — she can personalise",
        "Side menu has a red number showing how many are waiting",
    ],
    screenshot_label="/nurse Pending queue with 3 cards",
    notes_text=(
        "Open /nurse if possible. Show the actual queue.\n\n"
        "'This is what your nurse opens with her coffee in the morning. She works top to bottom. "
        "Each one takes maybe 30 seconds.'\n\n"
        "If they ask 'what if she's busy?' — 'The system holds pending requests for 24 hours. After "
        "that they auto-expire and the patient knows to submit again.'"
    ),
    page_num=10,
)

# =================================================================
# 11. WhatsApp magic
# =================================================================
feature_slide(
    section="The nurse",
    title="WhatsApp — without typing",
    sub="The clever bit. We don't pay Meta. Your nurse just sends from her own WhatsApp.",
    bullets=[
        "Tap 'Check with patient' → her WhatsApp opens with the message already typed",
        "She just hits Send — patient sees her real number, real face on display picture",
        "Five messages built in: Check · Confirm · Reschedule · Cancel · Reject · Reminder",
        "Every message uses your clinic's name and the patient's actual details automatically",
        "Zero per-message cost. Free forever.",
    ],
    screenshot_label="WhatsApp app opening with pre-typed message",
    notes_text=(
        "THIS is the slide that always lands hardest.\n\n"
        "'Most booking systems either don't do WhatsApp at all, OR they charge you for every "
        "message — 20 sen per send. 100 patients a week, 5 messages each, that's RM 400 a month "
        "just on messages.'\n\n"
        "'We use a trick called wa.me. The message is already written. She just taps Send. From her "
        "phone, her number, her account. Patient sees a normal WhatsApp conversation. Trust intact, "
        "zero ringgit per message.'\n\n"
        "Optionally show: 'And you can edit every template. If you want to write in Bahasa or Chinese, "
        "edit once and the change applies everywhere.'"
    ),
    page_num=11,
)

# =================================================================
# 12. Approve / track every message
# =================================================================
feature_slide(
    section="The nurse",
    title="Track every message — who sent what, when",
    sub="No more 'did anyone confirm Mrs Lim's appointment?' arguments.",
    bullets=[
        "Every WhatsApp message sent shows a green ✓ pill on the booking",
        "Hover to see: 'Confirmation sent by Nurse Aisha, 2 May 10:43 AM'",
        "All 5 message types tracked: Check, Confirm, Reject, Cancel, Reminder",
        "If two nurses share the shift, you instantly know which one did what",
        "Resend works the same way. Each send overwrites the latest record.",
    ],
    screenshot_label="Booking row with green ✓ Confirm and ✓ Reminder pills",
    notes_text=(
        "'No more guessing. No more \"I thought you sent it.\" Every send is logged.'\n\n"
        "Tie it back to a real scenario: 'Patient calls saying she never got a confirmation. You open "
        "the booking, see ✓ Confirm sent by Aisha at 10:43 AM yesterday. Done. No drama.'"
    ),
    page_num=12,
)

# =================================================================
# 13. Reminders
# =================================================================
feature_slide(
    section="The nurse",
    title="Reminders the day before",
    sub="The single biggest reason no-shows drop in clinics that adopt this.",
    bullets=[
        "Side menu: 'Send reminders' — shows tomorrow's confirmed appointments",
        "One tap per patient = WhatsApp opens with the reminder message ready",
        "She can do all 10–15 reminders in under 5 minutes the evening before",
        "Already-sent reminders show ✓ so she doesn't double-send",
        "Industry research: reminders cut no-shows by 30–50%",
    ],
    screenshot_label="Tomorrow's appointment list with 'Send reminder' buttons",
    notes_text=(
        "Math out the value: 'Say you get 2 no-shows a week. RM 150 each. That's RM 1,200 a month "
        "lost. Reminders typically halve that. So this feature alone saves you RM 600/month.'\n\n"
        "Then: 'The system costs RM 150 a month. Reminders alone pay for it 4 times over.'"
    ),
    page_num=13,
)

# =================================================================
# 14. Attendance
# =================================================================
feature_slide(
    section="The nurse",
    title="Mark who came in, who didn't",
    sub="Walk-in arrivals → search by IC → one tap.",
    bullets=[
        "'Today' filter in All Bookings shows just today's confirmed visits",
        "Patient walks in → search their IC → tap ✓ Attended",
        "Didn't show up? Tap 'No-show' instead — tracked separately",
        "Patient's visit count goes up automatically (used for 'returning patient' badges)",
        "Made a mistake? One-tap Undo. Everything reversible.",
    ],
    screenshot_label="Today filter with Attended / No-show buttons",
    notes_text=(
        "'Receptionist scans IC, types the last 4 digits, taps. Done.'\n\n"
        "Why no-show tracking matters: 'After a month you can see a pattern. Mrs Tan no-shows every "
        "other appointment. You decide if you require a deposit from her in future. Without this data "
        "you'd just be guessing.'"
    ),
    page_num=14,
)

# =================================================================
# 15. Book on behalf
# =================================================================
feature_slide(
    section="The nurse",
    title="When patients can't (or won't) use the form",
    sub="Elderly patients call. Walk-ins arrive. Auntie wants to book her sister.",
    bullets=[
        "'New booking' page — same form, but the nurse fills it for the patient",
        "Search existing patient by IC — name and phone auto-fill",
        "New patient? Fill once, saved forever.",
        "Same treatment + slot picker — works exactly like the public form",
        "Submits as 'confirmed' immediately — no waiting, no WhatsApp verification needed",
    ],
    screenshot_label="Staff new-booking page",
    notes_text=(
        "Anticipate the objection: 'But my elderly patients won't use the online form.'\n\n"
        "Answer: 'Of course — that's why your nurse has the same form on her side. Phone call comes "
        "in, she opens New Booking, types it in herself, done. The system doesn't require patients to "
        "do anything online. It just lets the ones who want to.'\n\n"
        "Strong point: 'Most of your appointments will still be made by your nurse. The system "
        "removes the FRICTION — phone tag, lost notes — not the human contact.'"
    ),
    page_num=15,
)


# =================================================================
# 16. DOCTOR DIVIDER
# =================================================================
section_divider("Part 3 of 4", "What your doctor sees",
                "Just their own day. Nothing extra.", 16)


# =================================================================
# 17. Doctor view
# =================================================================
feature_slide(
    section="The doctor",
    title="Their day, their patients, their schedule",
    sub="Doctors only see what's theirs. No nurse noise, no other doctors' calendars.",
    bullets=[
        "'Today' — list of patients coming in today, in time order, with treatment",
        "'My calendar' — week view of their own bookings only",
        "'My patients' — list of every patient they've ever treated",
        "'Block time' — mark lunch, MC days, conferences, weekends off",
        "Approved leaves automatically block their public bookings — no double-checking",
    ],
    screenshot_label="Doctor's Today view + My calendar",
    notes_text=(
        "Most clinic owners worry doctors will resist tech. Don't sell features here.\n\n"
        "'Doctors see their own calendar, their own patients. Nothing else. They can't see other "
        "doctors' bookings, can't see the owner's overview. Just their day. That's it.'\n\n"
        "If a doctor in the audience: speak directly: 'You'll know your day before you arrive. No "
        "surprises. Block your lunch once, the system protects it forever.'"
    ),
    page_num=17,
)


# =================================================================
# 18. OWNER DIVIDER
# =================================================================
section_divider("Part 4 of 4", "What you, the owner, see",
                "Everything. Edit anything. Hold everyone accountable.", 18)


# =================================================================
# 19. Owner dashboard + override
# =================================================================
feature_slide(
    section="The owner",
    title="You see everything, you can override anything",
    sub="Run a clinic with confidence — never blind to what's happening.",
    bullets=[
        "Overview dashboard: week's bookings, pending count, patient stats, repeat rate",
        "All bookings — filter by upcoming, today, attended, no-show, by doctor, by date",
        "Override any booking's status with one click (with optional note)",
        "Reset any staff password from the staff list",
        "Add or deactivate staff in 30 seconds",
    ],
    screenshot_label="/owner dashboard with stats",
    notes_text=(
        "You is plural here. Even if there are multiple owners (e.g., partners), each owner account "
        "sees the same.\n\n"
        "Real scenario to mention: 'Nurse accidentally rejects a booking she meant to approve. As "
        "owner, you can flip it back — and your override is logged in the audit trail so there's no "
        "blame game. You just fix it.'"
    ),
    page_num=19,
)


# =================================================================
# 20. Audit log + safety
# =================================================================
feature_slide(
    section="The owner",
    title="Every action recorded. Forever.",
    sub="Accountability built in. Required for Malaysian PDPA compliance.",
    bullets=[
        "Every approve, reject, cancel, override, message sent — recorded with timestamp and name",
        "Filter the log by action type, by staff member, by date",
        "Patient says 'I never agreed to this' — open the log and see exactly what happened",
        "Staff member leaves the clinic? Their name stays on every past record. Nothing erased.",
        "Required if a regulator asks 'who accessed patient X's data on date Y?'",
    ],
    screenshot_label="/owner/audit log with filters",
    notes_text=(
        "PDPA = Personal Data Protection Act. Malaysian law. Clinics handling patient data are "
        "subject to it.\n\n"
        "Don't dwell on legal — most clinic owners aren't fluent. Just frame it as 'protection': "
        "'If anything ever goes wrong — patient complaint, dispute, audit — you have a complete "
        "record of what happened. That's a lot harder when your records are in WhatsApp chats and "
        "notebooks.'"
    ),
    page_num=20,
)


# =================================================================
# 21. Team scheduling — duty + leave + custom shifts
# =================================================================
feature_slide(
    section="The owner",
    title="The whole team's schedule, on one page",
    sub="Duty calendar shows who's on, who's on leave, who's working half-day.",
    bullets=[
        "Default: everyone is on duty 9:00 AM – 9:00 PM, every day. No data entry needed.",
        "Doctor or nurse submits leave → you approve once → it auto-blocks their bookings",
        "Custom shift (e.g. half-day) → you approve → calendar reflects it",
        "Week view on phone, month view on desktop",
        "Working hours editor: change any doctor's weekly schedule in 1 minute",
    ],
    screenshot_label="Duty calendar — month view with leave + custom shift overlays",
    notes_text=(
        "Frame this as 'less work for you, the owner.'\n\n"
        "'You don't enter everyone's schedule. The default is everyone here, 9 to 9. Anyone going on "
        "leave just submits a request. You approve once. Done. The booking form automatically stops "
        "showing that doctor on those days. No manual coordination needed.'\n\n"
        "If they ask about complex shifts: 'You can configure a doctor to only work Tuesday and "
        "Thursday mornings. We change their working hours once. After that, the system honours it.'"
    ),
    page_num=21,
)


# =================================================================
# 22. WhatsApp templates editor
# =================================================================
feature_slide(
    section="The owner",
    title="All WhatsApp messages — edit once, applies everywhere",
    sub="You control the voice of your clinic. Change wording anytime.",
    bullets=[
        "All 6 templates editable: Check, Confirm Booking, Reschedule, Cancel, Reject, Reminder",
        "Live preview shows how the message will look with real patient data",
        "Tags like {patient_name}, {doctor_name}, {slot_label} auto-fill",
        "Add Bahasa Malaysia or Chinese version — change applies system-wide instantly",
        "Tone matches your clinic — formal, friendly, however you want",
    ],
    screenshot_label="Templates editor with side-by-side text + preview",
    notes_text=(
        "Quick demo of the editor if possible.\n\n"
        "'You decide whether messages start with \"Hi\" or \"Dear\" or \"Hello!\" — your clinic, your "
        "voice. The wording is yours. We just handle the plumbing.'"
    ),
    page_num=22,
)


# =================================================================
# 23. EXTRAS DIVIDER
# =================================================================
section_divider("The extras", "Things that make it feel finished",
                "Small touches your team will quietly appreciate.", 23)


# =================================================================
# 24. Mobile + home screen
# =================================================================
feature_slide(
    section="Bonus",
    title="Add to phone home screen — like a real app",
    sub="One tap to your clinic dashboard. No app store, no download, no waiting.",
    bullets=[
        "Open your dashboard on the phone browser, tap 'Add to Home Screen'",
        "Now there's a clinic icon on your home screen — looks like any other app",
        "One tap opens directly to today's bookings",
        "Works on iPhone and Android, same way",
        "No app store approval, no downloads, no updates to install",
    ],
    screenshot_label="iPhone home screen with the clinic app icon",
    notes_text=(
        "Quick demo on your own phone if you can: open the site, share, Add to Home Screen, point "
        "to the new icon.\n\n"
        "'No App Store delays. No 'please update the app' annoyances. It just works.'"
    ),
    page_num=24,
)


# =================================================================
# 25. Cost / pricing
# =================================================================
s = new_slide()
add_text(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6),
         "What it costs", size=28, bold=True, color=INK)
add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
         "Below the market. Way below.", size=14, color=MUTED)

# Comparison table
y = Inches(2.1)
rows = [
    ("Calendly (generic)", "RM 65 / mo", "No WhatsApp, no clinic features, no roles", MUTED),
    ("SimplyBook.me (closest rival)", "RM 55 / mo", "Generic global tool, not clinic-shaped", MUTED),
    ("Aoikumo (Malaysian, full-suite)", "RM 223–557 / mo", "Forces full migration, includes EMR/billing", MUTED),
    ("Our system, retail", "RM 150 / mo", "Branded as your clinic, built for Malaysian dental", BRAND_DARK),
    ("Our system, friend rate", "RM 80 / mo", "Same product, smaller bill, first 3 months free", BRAND),
]
col_widths = [Inches(4.5), Inches(2.5), Inches(5.5)]
header_y = y
for i, w in enumerate(col_widths):
    x = Inches(0.8) + sum(col_widths[:i], Inches(0))
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, header_y, w, Inches(0.45))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BRAND_DARK
    rect.line.fill.background()
    label = ["Option", "Price", "Trade-off"][i]
    add_text(s, x + Inches(0.15), header_y + Inches(0.08), w, Inches(0.3),
             label, size=12, bold=True, color=WHITE)

for ridx, (name, price, note, color) in enumerate(rows):
    ry = header_y + Inches(0.45) + Inches(0.55) * ridx
    bg_color = WHITE if ridx < 3 else BRAND_LIGHT
    for i, w in enumerate(col_widths):
        x = Inches(0.8) + sum(col_widths[:i], Inches(0))
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ry, w, Inches(0.55))
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        cell.line.color.rgb = MUTED
        cell.line.width = Pt(0.25)
        val = [name, price, note][i]
        bold = ridx >= 3 and i < 2
        add_text(s, x + Inches(0.15), ry + Inches(0.13), w, Inches(0.35),
                 val, size=12, bold=bold, color=color if i < 2 else INK_SOFT)

# One-time setup callout
y = Inches(5.6)
rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y,
                          Inches(12.5), Inches(1.0))
rect.fill.solid()
rect.fill.fore_color.rgb = WHITE
rect.line.color.rgb = BRAND
rect.line.width = Pt(1.5)
add_text(s, Inches(1.0), y + Inches(0.15), Inches(11), Inches(0.4),
         "One-time setup: RM 1,000", size=15, bold=True, color=BRAND_DARK)
add_text(s, Inches(1.0), y + Inches(0.55), Inches(11), Inches(0.4),
         "Includes onboarding your staff, importing existing patients, customising templates, "
         "and training. Free fixes for the first 60 days — money-back guarantee if it doesn't help.",
         size=11, color=INK_SOFT)

add_footer(s, SW, SH, 25, TOTAL)
add_notes(s,
    "Pricing is touchy with friends. Lead with comparison, not your number.\n\n"
    "'Aoikumo would charge them RM 250+ a month — and require them to migrate everything. Calendly "
    "is RM 65 but is built for individual freelancers, not a clinic. We sit right in the middle: "
    "clinic-shaped, Malaysian-shaped, RM 150 retail.'\n\n"
    "Then drop the friend rate: 'For you, RM 80 a month. First 3 months free. One-time RM 1,000 to "
    "cover setup. That's it. Cancel any month.'\n\n"
    "If they want lower: 'I can drop the setup to RM 500 if cash is tight. The RM 80 monthly stays.'"
)


# =================================================================
# 26. Risk reversal + close
# =================================================================
s = new_slide()
add_text(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6),
         "What if it doesn't work for you?", size=28, bold=True, color=INK)
add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
         "No commitments. No vendor lock-in. Walk away anytime.",
         size=14, color=MUTED)

items = [
    ("60-day money-back",
     "Setup fee is fully refunded if you decide it doesn't help in the first 60 days. No questions asked."),
    ("Cancel anytime",
     "Monthly billing. Stop the month you want. No contracts, no termination fees."),
    ("Your data is yours",
     "We export your patient list and bookings as CSV any time you ask. Take it elsewhere — or print it. It's yours."),
    ("It works alongside what you have",
     "Doesn't replace your EMR or accounting. Use it for bookings only. Keep everything else."),
]
y = Inches(2.4)
card_w = Inches(2.95)
gap = Inches(0.15)
total_w = card_w * 4 + gap * 3
start_x = (SW - total_w) / 2
for i, (title, body) in enumerate(items):
    x = start_x + (card_w + gap) * i
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(3.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BRAND
    rect.line.width = Pt(1.5)
    cir = s.shapes.add_shape(MSO_SHAPE.OVAL, x + card_w/2 - Inches(0.4),
                             y + Inches(0.3), Inches(0.8), Inches(0.8))
    cir.fill.solid()
    cir.fill.fore_color.rgb = GREEN_BG
    cir.line.fill.background()
    add_text(s, x, y + Inches(0.45), card_w, Inches(0.5),
             "✓", size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.4),
             Inches(0.5), title, size=14, bold=True, color=BRAND_DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(2.0), card_w - Inches(0.4),
             Inches(1.1), body, size=11, color=INK_SOFT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
         "The decision isn't 'do I trust this system forever'.",
         size=16, color=INK_SOFT, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6),
         "It's 'do I want to try it for 60 days, knowing I can walk away with my data intact?'",
         size=16, color=BRAND_DARK, bold=True, align=PP_ALIGN.CENTER)

add_footer(s, SW, SH, 26, TOTAL)
add_notes(s,
    "This slide does all the closing work for you. Don't oversell — just read the cards.\n\n"
    "'You can walk away with your data. You can cancel any month. The first 60 days are guaranteed. "
    "What's stopping us from trying?'\n\n"
    "Then shut up. Let them respond. The longer the pause, the better."
)


# =================================================================
# 27. Q&A
# =================================================================
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.18), SW, SH - Inches(0.18))
bg.fill.solid()
bg.fill.fore_color.rgb = PAPER
bg.line.fill.background()

add_text(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.7),
         "Questions?", size=64, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.5),
         "What would worry you about using this on Monday morning?",
         size=20, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

add_text(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
         "Or: tell me one thing you wish a clinic system did",
         size=14, color=BRAND_DARK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.4),
         "that none of them currently do.",
         size=14, color=BRAND_DARK, align=PP_ALIGN.CENTER)

add_footer(s, SW, SH, 27, TOTAL)
add_notes(s,
    "Don't fill the silence. The pause IS the close.\n\n"
    "Common questions and quick answers:\n\n"
    "Q: 'Internet down — what happens?'\n"
    "A: 'You go back to manual for a few hours. Once internet's back, the system catches up. We have "
    "Supabase's regional backups so no data is ever lost — even if our server goes down.'\n\n"
    "Q: 'How long until it's up?'\n"
    "A: 'Could deploy your clinic today. Realistic: 3–5 days to onboard staff, import patients, "
    "customise templates.'\n\n"
    "Q: 'What if you stop building it?'\n"
    "A: 'Your data exports anytime. The source code is documented. Worst case, you migrate to "
    "Calendly with your patient list intact. You're never trapped.'\n\n"
    "Q: 'Can my accountant see the data?'\n"
    "A: 'You can give them an owner login OR export reports as CSV. Both work.'\n\n"
    "Q: 'Multiple branches?'\n"
    "A: 'Yes — that's actually next on the roadmap. Right now, one clinic per setup. Multi-branch "
    "coming in a few months.'"
)


# =================================================================
# Save
# =================================================================
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "Goodcare_Marketing_Deck.pptx")
prs.save(out_path)
print(f"Wrote {out_path}")
print(f"Slides: {len(prs.slides)}")
