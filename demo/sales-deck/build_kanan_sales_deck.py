"""Generate the Kanan clinic-booking sales deck.

Outputs ./demo/sales-deck/Kanan_Clinic_Booking_Sales_Deck.pptx — 18 slides,
Kanan brand styling (navy + gold + warm-white per the brand guidelines memory).
Content mirrors ./kanan-sales-deck.md — edit the .md for copy, then re-run.

Slides 9–12 are NEW (v2): shared-kiosk lockscreen, Premium room flow,
staff cards on /owner/staff, and Approvals + Audit log. All screenshot
boxes are LABELLED PLACEHOLDERS — capture from clinic-booking-kanan-team.vercel.app
and drop your real PNGs into the corresponding rectangles in PowerPoint.

Run:  python3 demo/sales-deck/build_kanan_sales_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Kanan brand colors (must match brand guidelines memory) ─────────────
NAVY        = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_LIGHT  = RGBColor(0x2E, 0x43, 0x74)
GOLD        = RGBColor(0xC9, 0xA2, 0x27)
GOLD_LIGHT  = RGBColor(0xE3, 0xC7, 0x6A)
WARM_WHITE  = RGBColor(0xF4, 0xF1, 0xEA)
GREY        = RGBColor(0x6B, 0x72, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x05, 0x96, 0x69)
RED         = RGBColor(0xC0, 0x39, 0x2B)

# Font choice: Helvetica is safe across MY-installed machines.
# Once a customer opens the file, the brand guideline fonts (Fraunces /
# Hanken Grotesk) need to be installed locally to actually render.
HEAD_FONT = "Helvetica"
BODY_FONT = "Helvetica"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)


# ─── Helpers ──────────────────────────────────────────────────────────────

def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(shape, color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, italic=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=NAVY, bullet_color=GOLD,
                line_gap_pt=6, bullet_char="●"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(line_gap_pt)
        b = p.add_run()
        b.text = f"{bullet_char}  "
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        b.font.name = BODY_FONT
        t = p.add_run()
        t.text = item
        t.font.size = Pt(size)
        t.font.color.rgb = color
        t.font.name = BODY_FONT
    return box


def add_speaker_note(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_footer(slide, page_num, total):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.4),
             Inches(6), Inches(0.3),
             "Kanan Digital Enterprise · kanan.my",
             size=10, color=GREY, font=BODY_FONT)
    add_text(slide, SLIDE_W - MARGIN - Inches(1), SLIDE_H - Inches(0.4),
             Inches(1), Inches(0.3), f"{page_num} / {total}",
             size=10, color=GREY, align=PP_ALIGN.RIGHT, font=BODY_FONT)


def add_screenshot_placeholder(slide, left, top, width, height, caption=""):
    box = add_rect(slide, left, top, width, height, WARM_WHITE, line_color=GOLD_LIGHT)
    box.line.width = Pt(1.5)
    add_text(slide, left, top, width, height, "[ screenshot ]",
             size=11, italic=True, color=GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=BODY_FONT)
    if caption:
        add_text(slide, left, top + height + Inches(0.05), width, Inches(0.3),
                 caption, size=10, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER, font=BODY_FONT)


# ─── Slide builders ───────────────────────────────────────────────────────

def s1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WARM_WHITE)
    # Gold accent strip top
    add_rect(s, 0, 0, SLIDE_W, Inches(0.15), GOLD)
    # Centered wordmark area (image goes here — placeholder until pptx supports embedding)
    add_text(s, 0, Inches(2.2), SLIDE_W, Inches(0.6),
             "kanan",
             size=80, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             font="Georgia")  # serif stand-in for the actual logo
    add_text(s, 0, Inches(3.2), SLIDE_W, Inches(0.5),
             "your trusted right hand",
             size=18, italic=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.1), SLIDE_W, Inches(0.6),
             "A booking system built for Malaysian dental clinics",
             size=22, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.8), SLIDE_W, Inches(0.4),
             "by Kanan Digital Enterprise",
             size=14, color=GREY, align=PP_ALIGN.CENTER)
    # Bottom strip
    add_rect(s, 0, SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), NAVY)
    add_text(s, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.3),
             "kanan.my · hello@kanan.my", size=11, color=WARM_WHITE,
             align=PP_ALIGN.CENTER)


def slide_title_only(prs, page, total, title, lead=None):
    """Common header for body slides — title bar + optional lead text."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WARM_WHITE)
    # Gold thin accent
    add_rect(s, MARGIN, Inches(0.55), Inches(0.5), Inches(0.05), GOLD)
    add_text(s, MARGIN, Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.8),
             title, size=32, bold=True, color=NAVY, font=HEAD_FONT)
    if lead:
        add_text(s, MARGIN, Inches(1.6), SLIDE_W - 2 * MARGIN, Inches(0.6),
                 lead, size=16, italic=True, color=NAVY_LIGHT)
    add_footer(s, page, total)
    return s


def s2_pain(prs, total):
    s = slide_title_only(prs, 2, total, "Short-staffed and stretched thin?",
                         "Most days end the same way at a small dental clinic:")
    add_bullets(s, MARGIN, Inches(2.4), SLIDE_W - 2 * MARGIN, Inches(4), [
        "50+ WhatsApp messages a day — just to confirm and remind appointments",
        "A no-show wastes a whole 1-hour slot, and the patient never paid a deposit",
        "The owner has no idea who's actually busy and who's coasting",
        "Patient records live half in a notebook, half in a Google Sheet, half in someone's head",
        "When a doctor or nurse takes leave, the schedule unravels for a week",
    ], size=18, line_gap_pt=10)
    add_speaker_note(s, "Let them nod. Don't oversell. Ask: \"Which of these sounds like your Tuesday?\"")


def s3_meet_kanan(prs, total):
    s = slide_title_only(prs, 3, total,
                         "We've been there. So we built the right hand.")
    add_text(s, MARGIN, Inches(1.8), SLIDE_W - 2 * MARGIN, Inches(0.9),
             "Kanan (\"KAH-nahn\") means right in Malay — as in right hand. "
             "The dependable hand on the right side of your business.",
             size=16, color=NAVY_LIGHT, italic=True)
    # Two founder cards
    card_w = Inches(5.5)
    card_h = Inches(2.2)
    card_top = Inches(3.2)
    gap = Inches(0.3)
    total_w = card_w * 2 + gap
    left_start = (SLIDE_W - total_w) // 2
    for i, (name, role) in enumerate([
        ("[Founder name 1]", "Co-founder · [role]"),
        ("[Founder name 2]", "Co-founder · [role]"),
    ]):
        l = left_start + i * (card_w + gap)
        add_rect(s, l, card_top, card_w, card_h, WHITE, line_color=NAVY)
        add_text(s, l, card_top + Inches(0.4), card_w, Inches(0.5),
                 "[photo]", size=12, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER)
        add_text(s, l, card_top + Inches(1.1), card_w, Inches(0.5),
                 name, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, l, card_top + Inches(1.6), card_w, Inches(0.4),
                 role, size=13, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(5.7), SLIDE_W - 2 * MARGIN, Inches(0.5),
             "A Malaysian software studio. We build for the businesses around us — "
             "not for venture capital.",
             size=14, color=NAVY, align=PP_ALIGN.CENTER, italic=True)
    add_speaker_note(s, "Personalize. Use first names. Don't say 'we leverage AI' — say 'we built a thing that does X'.")


def s4_product_glance(prs, total):
    s = slide_title_only(prs, 4, total, "One product. The whole booking lifecycle.")
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) // 3
    cols = [
        ("📱", "Patient self-books",
         "via your /book page. Language toggle EN · 中文 · BM."),
        ("💬", "Nurse runs the day",
         "Pending queue, WhatsApp confirmation, reminders, recall."),
        ("📊", "Owner sees everything",
         "Bookings, doctor + nurse performance, audit log."),
    ]
    for i, (icon, h, body) in enumerate(cols):
        l = MARGIN + i * (col_w + Inches(0.3))
        add_rect(s, l, Inches(2.3), col_w, Inches(2.5), WHITE, line_color=NAVY_LIGHT)
        add_text(s, l, Inches(2.45), col_w, Inches(0.6), icon,
                 size=36, align=PP_ALIGN.CENTER)
        add_text(s, l, Inches(3.2), col_w, Inches(0.5),
                 h, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, l, Inches(3.8), col_w, Inches(0.8),
                 body, size=12, color=GREY, align=PP_ALIGN.CENTER)
    add_screenshot_placeholder(s, MARGIN, Inches(5.1),
                               SLIDE_W - 2 * MARGIN, Inches(1.6),
                               "Screenshot collage: /book mobile · /nurse desktop · /owner dashboard")
    add_speaker_note(s, "30-second product. Three audiences, three views.")


def s5_patient(prs, total):
    s = slide_title_only(prs, 5, total,
                         "Bookings without phone tag.",
                         "Your patient opens your /book page on their phone. 4 taps to a request.")
    add_bullets(s, MARGIN, Inches(2.6), Inches(7), Inches(4.5), [
        "Identity first (IC/passport + WhatsApp), THEN pick what they want — book / reschedule / cancel",
        "Language toggle: EN · 中文 · BM (name input stays Romanized for IC)",
        "5-week future window. Real slot availability — no requests for slots that aren't free",
        "Submitted requests land in your nurse's queue, never directly on the calendar",
    ], size=15, line_gap_pt=10)
    add_screenshot_placeholder(s, Inches(8), Inches(2.5), Inches(4.7), Inches(4),
                               "/book mobile screenshot")
    add_speaker_note(s, "Show the live demo here if there's bandwidth.")


def s6_nurse(prs, total):
    s = slide_title_only(prs, 6, total,
                         "Your nurse stops being a router.",
                         "Every patient is one click away from confirmed, with the WhatsApp message pre-written.")
    add_bullets(s, MARGIN, Inches(2.6), Inches(7), Inches(4.5), [
        "6 WhatsApp templates editable by owner — check / confirm / cancel / reject / reminder / recall",
        "Click 'Send' → WhatsApp opens with the message + slot pre-filled, zero retyping",
        "Recall worklist surfaces patients overdue for 6-month checkup automatically",
        "Every action audit-logged with the nurse's name + timestamp",
    ], size=15, line_gap_pt=10)
    add_screenshot_placeholder(s, Inches(8), Inches(2.5), Inches(4.7), Inches(4),
                               "/nurse pending queue screenshot")
    add_speaker_note(s, "Highlight the audit log angle — owners love that they can verify.")


def s7_doctor(prs, total):
    s = slide_title_only(prs, 7, total,
                         "Your doctor sees the day, marks attendance, that's it.")
    add_bullets(s, MARGIN, Inches(2.4), Inches(7), Inches(4.5), [
        "Clean today view of their bookings, on a 9-9 grid",
        "One-tap Attended or No-show after each patient",
        "Block time (lunch / surgery prep / personal) without bothering the nurse",
        "Their own utilization heatmap on Premium",
    ], size=15, line_gap_pt=10)
    add_screenshot_placeholder(s, Inches(8), Inches(2.3), Inches(4.7), Inches(4.2),
                               "/doctor today view screenshot")
    add_speaker_note(s, "Doctors get value too — less admin interruptions during chair time.")


def s8_owner_dash(prs, total):
    s = slide_title_only(prs, 8, total, "Finally, you can see the clinic.")
    add_bullets(s, MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(4.5), [
        "This week's bookings + repeat rate + pending count — at a glance",
        "Doctor performance (Premium): bookings, attended, no-show, attendance % — per doctor",
        "Nurse performance (Premium): bookings created, approvals, reminders + recalls, attendance marks",
        "Chair utilization heatmap (Premium): which chair-hours are actually filled",
        "Full audit log — who did what, when, with before/after snapshots",
    ], size=15, line_gap_pt=10)
    add_screenshot_placeholder(s, MARGIN, Inches(5.0), SLIDE_W - 2 * MARGIN, Inches(1.7),
                               "Screenshot grid: doctor-perf + nurse-perf + utilization")
    add_speaker_note(s, "This is where upgrade pressure to Premium happens organically.")


def tier_card(slide, left, top, width, height, tier, price, seats, features, accent=NAVY):
    add_rect(slide, left, top, width, height, WHITE, line_color=accent)
    # Top color band
    add_rect(slide, left, top, width, Inches(0.15), accent)
    add_text(slide, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.6),
             tier, size=22, bold=True, color=NAVY)
    add_text(slide, left + Inches(0.3), top + Inches(0.95), width - Inches(0.6), Inches(0.5),
             price, size=18, bold=True, color=GOLD)
    add_text(slide, left + Inches(0.3), top + Inches(1.5), width - Inches(0.6), Inches(0.4),
             seats, size=12, color=GREY, italic=True)
    add_bullets(slide, left + Inches(0.3), top + Inches(2.0),
                width - Inches(0.6), height - Inches(2.2),
                features, size=11, line_gap_pt=5, bullet_char="✓",
                bullet_color=GREEN)


def s9_terminal(prs, total):
    s = slide_title_only(prs, 9, total,
                         "One tablet at the front desk runs the whole day.",
                         "Plug in a cheap Android or iPad at reception — anyone uses it, every action is PIN-locked.")
    add_bullets(s, MARGIN, Inches(2.4), Inches(7), Inches(4.5), [
        "Big clock + branded background (matches your theme)",
        "Upcoming patients panel — next 48 hours, time + name + IC + doctor",
        "Live count tiles: Pending · Recalls · Today · Reminders",
        "Mark attended / no-show inline — one PIN, one tap",
        "Marked patients collapse under a 'Marked · 5 attended · 1 no-show' bar",
        "Doctors can also PIN in to view their own day + check-out + submit HR",
    ], size=14, line_gap_pt=8)
    add_screenshot_placeholder(s, Inches(7.9), Inches(2.3), Inches(4.8), Inches(4.4),
                               "/home — lockscreen with upcoming-patients panel")
    add_speaker_note(s,
        "Reception runs on one device. Nobody shares passwords — the kiosk is signed in as 'terminal' "
        "and every booking-changing action requires a 6-digit PIN from a real nurse or doctor. "
        "Capture: open /home on a touchscreen and tap Pending to show the PIN modal.")


def s10_room_flow(prs, total):
    s = slide_title_only(prs, 10, total,
                         "Premium room flow — patient arrives → treatment done.",
                         "Three taps, two PINs, zero paper. Every step audit-logged.")
    # Three-step card row
    card_w = Inches(3.95)
    card_h = Inches(2.6)
    card_top = Inches(2.4)
    gap = Inches(0.15)
    total_w = card_w * 3 + gap * 2
    left_start = (SLIDE_W - total_w) // 2
    steps = [
        ("1", "Patient arrives", "Nurse taps Check in → picks room → Nurse PIN → ✓ In Room 2", NAVY_LIGHT),
        ("2", "Doctor's turn", "Doctor PIN at lockscreen → sees 'In Room 2' badge on their booking row", GOLD),
        ("3", "After the visit", "Doctor taps Check out → picks treatment done → Doctor PIN → ✓ Done · Crown fitting", GREEN),
    ]
    for i, (n, title, body, accent) in enumerate(steps):
        x = left_start + i * (card_w + gap)
        add_rect(s, x, card_top, card_w, card_h, WHITE, line_color=accent)
        add_rect(s, x, card_top, card_w, Inches(0.15), accent)
        add_text(s, x + Inches(0.25), card_top + Inches(0.35), card_w - Inches(0.5), Inches(0.5),
                 n, size=28, bold=True, color=NAVY, font="Georgia")
        add_text(s, x + Inches(0.25), card_top + Inches(0.95), card_w - Inches(0.5), Inches(0.4),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), card_top + Inches(1.4), card_w - Inches(0.5), Inches(1.1),
                 body, size=11, color=GREY)
    # Detail strip below
    add_bullets(s, MARGIN, Inches(5.3), SLIDE_W - 2 * MARGIN, Inches(1.5), [
        "Rooms list editable by owner — 'Surgery', 'Hygiene bay 2', whatever your clinic uses",
        "Treatment-done text travels into the visit history + audit log automatically",
        "Auto-marks attended + bumps visit_count + resets the recall timer on check-out",
    ], size=12, line_gap_pt=5)
    add_screenshot_placeholder(s, Inches(9.5), Inches(5.3), Inches(3.2), Inches(1.5),
                               "Optional: room picker / treatment picker modal")
    add_speaker_note(s,
        "Walk through with finger on screen: 'Nurse — doctor — done.' Sell the visibility: "
        "owner can see at a glance which rooms are busy. Capture: open the lockscreen + tap "
        "Check in on a pending row to show the room picker.")


def s11_staff_cards(prs, total):
    s = slide_title_only(prs, 11, total,
                         "Manage your team from one place.",
                         "Doctors and nurses — one card each. Expand, edit, save inline.")
    add_bullets(s, MARGIN, Inches(2.3), Inches(7), Inches(4.5), [
        "Working hours — per-weekday editor (both doctors AND nurses)",
        "Leave entitlement — Annual / MC / Emergency days, owner sets per person",
        "Doctor profile (Premium) — expertise tags + bio + ★ rating for /book cards",
        "Recent activity — last 5 leave + shift events with status pills",
        "Account actions — PIN reset / password reset / deactivate",
        "Cards collapsed by default — header shows '5 days/wk · custom' summary",
    ], size=13, line_gap_pt=8)
    add_screenshot_placeholder(s, Inches(7.9), Inches(2.3), Inches(4.8), Inches(4.4),
                               "/owner/staff — one card expanded, one collapsed")
    add_speaker_note(s,
        "Before this redesign, owners had to click 'Working hours', 'My account', 'Doctors & nurses' — "
        "three pages. Now it's one card per person. Capture: open /owner/staff and click a doctor "
        "card to expand the Working hours + Leave entitlement sections.")


def s12_approvals_audit(prs, total):
    s = slide_title_only(prs, 12, total,
                         "Compliance without the paperwork.",
                         "Every HR decision and every booking action — on the record.")
    # Two-column layout
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) // 2
    col_top = Inches(2.3)
    col_h = Inches(4.4)
    # Left col — Approvals
    add_rect(s, MARGIN, col_top, col_w, col_h, WHITE, line_color=NAVY_LIGHT)
    add_rect(s, MARGIN, col_top, col_w, Inches(0.15), NAVY_LIGHT)
    add_text(s, MARGIN + Inches(0.3), col_top + Inches(0.3), col_w - Inches(0.6), Inches(0.5),
             "Approvals", size=18, bold=True, color=NAVY)
    add_bullets(s, MARGIN + Inches(0.3), col_top + Inches(0.9),
                col_w - Inches(0.6), col_h - Inches(1.2),
                [
                    "Pending leave + pending shift changes side-by-side",
                    "One-tap Approve / Reject with optional reviewer note",
                    "Approved leave auto-blocks the doctor calendar",
                    "Permanent shift changes auto-update working hours on approval",
                ], size=12, line_gap_pt=6)
    # Right col — Audit log
    right_left = MARGIN + col_w + Inches(0.3)
    add_rect(s, right_left, col_top, col_w, col_h, WHITE, line_color=GOLD)
    add_rect(s, right_left, col_top, col_w, Inches(0.15), GOLD)
    add_text(s, right_left + Inches(0.3), col_top + Inches(0.3), col_w - Inches(0.6), Inches(0.5),
             "Audit log (Premium)", size=18, bold=True, color=NAVY)
    add_bullets(s, right_left + Inches(0.3), col_top + Inches(0.9),
                col_w - Inches(0.6), col_h - Inches(1.2),
                [
                    "Last 500 actions — actor, time, patient, 'X → Y' transition",
                    "Filter by action / actor / entity + free-text search",
                    "'Show raw JSON' expander for the full payload",
                    "Covers: bookings, PINs, leave, theme, rooms, WhatsApp sends",
                ], size=12, line_gap_pt=6)
    add_screenshot_placeholder(s, MARGIN, Inches(6.95), SLIDE_W - 2 * MARGIN, Inches(0.3),
                               "Optional: /owner/hr-approvals + /owner/audit screenshots")
    add_speaker_note(s,
        "PDPA + SSM both ask 'who did this and when?' — your answer is one filter + a screenshot. "
        "Lean on this slide for multi-doctor clinics worried about discipline or disputes. "
        "Capture: /owner/hr-approvals with at least one pending row visible, and /owner/audit "
        "showing a few booking rows with patient names + status transitions.")


def s13_standard(prs, total):
    s = slide_title_only(prs, 13, total, "Standard — RM 150 / month",
                         "Everything you need to run a single dental clinic that's drowning in WhatsApp.")
    tier_card(s, MARGIN, Inches(2.4), SLIDE_W - 2 * MARGIN, Inches(4.4),
              "Standard tier", "RM 150 / month",
              "1 owner  ·  2 doctors  ·  3 nurses (top-up on WhatsApp)",
              [
                  "Patient self-book + reschedule + cancel  (with EN / CH / BM toggle)",
                  "Nurse pending queue + 6 editable WhatsApp templates",
                  "Recall reminders (every 6 months, configurable)",
                  "Duty calendar — doctor schedules",
                  "Clinical calendar + per-doctor working hours editor",
                  "Branding (your logo, brand color, font, button shape)",
                  "Daily / weekly / monthly CSV email backup",
                  "Overview dashboard (basic) + mobile-friendly throughout",
              ], accent=NAVY_LIGHT)
    add_text(s, MARGIN, Inches(6.95), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Need more seats? WhatsApp us — we top up without forcing an upgrade.",
             size=11, italic=True, color=GREY)


def s14_premium(prs, total):
    s = slide_title_only(prs, 14, total, "Premium — RM 280 / month",
                         "Everything in Standard, plus the visibility you stop being able to live without.")
    tier_card(s, MARGIN, Inches(2.4), SLIDE_W - 2 * MARGIN, Inches(4.4),
              "Premium tier", "RM 280 / month",
              "1 owner  ·  4 doctors  ·  6 nurses (top-up on WhatsApp)",
              [
                  "Room flow — nurse check-in to room, doctor check-out + treatment-done",
                  "Patient-facing doctor cards on /book — expertise + ★ rating",
                  "Rooms editor — owner names the operatories your clinic uses",
                  "Doctor + nurse performance dashboards — attendance, no-shows, activity",
                  "Chair utilization heatmap — which chair-hours are actually filled",
                  "Duty calendar includes nurses, not just doctors",
                  "Audit log — patient name + 'X → Y' transitions, search + filter",
                  "Roadmap: review system + payroll handoff (PayrollPanda)",
              ], accent=GOLD)
    add_text(s, MARGIN, Inches(6.95), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Seat cap reached? WhatsApp us — we top up without forcing a Franchise upgrade.",
             size=11, italic=True, color=GREY)


def s15_comparison(prs, total):
    s = slide_title_only(prs, 15, total, "Standard vs Premium at a glance")
    headers = ["", "Standard", "Premium"]
    rows = [
        ("Monthly price", "RM 150", "RM 280"),
        ("Owner / Doctor / Nurse seats", "1 / 2 / 3", "1 / 4 / 6"),
        ("Booking + WhatsApp templates", "✓", "✓"),
        ("Recall reminders", "✓", "✓"),
        ("Terminal lockscreen + nurse PIN flow", "✓", "✓"),
        ("Staff cards (collapsible per-person editor)", "✓", "✓"),
        ("Approvals (leave + shift in one place)", "✓", "✓"),
        ("CSV backup (folder + auto-email)", "✓", "✓"),
        ("Room flow — check-in / out / treatment", "—", "✓"),
        ("Patient-facing doctor cards on /book", "—", "✓"),
        ("Duty calendar — nurses", "—", "✓"),
        ("Doctor + nurse performance analytics", "—", "✓"),
        ("Chair utilization heatmap", "—", "✓"),
        ("Audit log (full PII + transitions)", "—", "✓"),
    ]
    table_left = MARGIN
    table_top = Inches(2.0)
    col_w = [Inches(5.5), Inches(3.2), Inches(3.2)]
    # Tighter rows than v1 — table grew from 11 → 14 rows to cover new features.
    row_h = Inches(0.32)
    # Header
    for i, h in enumerate(headers):
        x = table_left + sum(col_w[:i], Inches(0))
        add_rect(s, x, table_top, col_w[i], row_h, NAVY)
        add_text(s, x, table_top, col_w[i], row_h, h,
                 size=13, bold=True, color=WARM_WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE,
                 font=HEAD_FONT)
    # Rows
    for ri, row in enumerate(rows):
        y = table_top + row_h * (ri + 1)
        bg = WARM_WHITE if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            x = table_left + sum(col_w[:ci], Inches(0))
            add_rect(s, x, y, col_w[ci], row_h, bg)
            is_check = val == "✓"
            is_dash = val == "—"
            color = GREEN if is_check else (GREY if is_dash else NAVY)
            add_text(s, x + Inches(0.15), y, col_w[ci] - Inches(0.3), row_h, val,
                     size=12,
                     bold=ci == 0 or is_check,
                     color=color,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
    add_speaker_note(s, "Let them stare at this for 10 seconds. Don't talk through it line by line.")


def s16_how_we_work(prs, total):
    s = slide_title_only(prs, 16, total, "Start to finish — no surprises")
    steps = [
        ("1", "Talk to our team", "WhatsApp or quick call. Tell us what's slowing you down."),
        ("2", "Schedule a discussion", "Longer working session. No sales deck."),
        ("3", "Workflow & framework", "Map your current ops + where software fits."),
        ("4", "Proposal & strategic review", "1-page scope, timeline, price. Push back where it doesn't fit."),
        ("5", "Demo", "Working prototype on your real data, in weeks not months."),
        ("6", "Commercial discussion", "Full commercial agreed once the demo proves out."),
        ("7", "Setup & training", "Go-live, team trained, real users on the system. WhatsApp throughout."),
        ("8", "Done — and still here", "Support continues. We don't disappear after handover."),
    ]
    grid_top = Inches(2.0)
    cell_w = (SLIDE_W - 2 * MARGIN - Inches(0.3) * 3) // 4
    cell_h = Inches(2.1)
    for i, (n, title, body) in enumerate(steps):
        row, col = divmod(i, 4)
        x = MARGIN + col * (cell_w + Inches(0.3))
        y = grid_top + row * (cell_h + Inches(0.3))
        add_rect(s, x, y, cell_w, cell_h, WHITE, line_color=NAVY_LIGHT)
        # Gold accent line
        add_rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.04), GOLD)
        add_text(s, x + Inches(0.25), y + Inches(0.4), cell_w - Inches(0.5), Inches(0.5),
                 n, size=28, bold=True, color=NAVY, font="Georgia")
        add_text(s, x + Inches(0.25), y + Inches(0.95), cell_w - Inches(0.5), Inches(0.45),
                 title, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), y + Inches(1.4), cell_w - Inches(0.5), Inches(0.7),
                 body, size=10, color=GREY)
    add_speaker_note(s, "Step 8 is the differentiator. Lead with it: 'We don't ghost you after launch.'")


def s17_trust(prs, total):
    s = slide_title_only(prs, 17, total, "Your data stays yours",
                         "Trust before technology.")
    add_bullets(s, MARGIN, Inches(2.4), SLIDE_W - 2 * MARGIN, Inches(4), [
        "Hosted in Malaysia wherever possible (Supabase + Vercel)",
        "PDPA-aware — patient data export any time, schema visible on request",
        "We don't train external AI models on your operational data",
        "Owner downloads a full CSV backup any day. Daily auto-email backup if you want it.",
        "Sterilization log, X-rays, clinical treatment notes stay in your EMR — not in our system",
    ], size=15, line_gap_pt=10)
    add_text(s, MARGIN, Inches(6.4), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Kanan Digital Enterprise (SSM Enterprise registration in progress) · Kuala Lumpur, Malaysia",
             size=11, italic=True, color=GREY)


def s18_try_it(prs, total):
    s = slide_title_only(prs, 18, total, "Try the demos. Get set up in 4 weeks.")
    # Two demo cards
    card_w = Inches(5.5)
    card_h = Inches(2.3)
    card_top = Inches(2.2)
    gap = Inches(0.3)
    left_start = (SLIDE_W - (card_w * 2 + gap)) // 2
    # Standard demo card
    add_rect(s, left_start, card_top, card_w, card_h, WHITE, line_color=NAVY_LIGHT)
    add_rect(s, left_start, card_top, card_w, Inches(0.15), NAVY_LIGHT)
    add_text(s, left_start, card_top + Inches(0.4), card_w, Inches(0.5),
             "Standard demo", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, left_start, card_top + Inches(1.0), card_w, Inches(0.5),
             "standard-demo.kanan.my", size=14, color=GOLD,
             align=PP_ALIGN.CENTER, bold=True)
    add_text(s, left_start, card_top + Inches(1.5), card_w, Inches(0.6),
             "RM 150/month · 2 doctors · 3 nurses",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # Premium demo card
    l2 = left_start + card_w + gap
    add_rect(s, l2, card_top, card_w, card_h, WHITE, line_color=GOLD)
    add_rect(s, l2, card_top, card_w, Inches(0.15), GOLD)
    add_text(s, l2, card_top + Inches(0.4), card_w, Inches(0.5),
             "Premium demo", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, l2, card_top + Inches(1.0), card_w, Inches(0.5),
             "premium-demo.kanan.my", size=14, color=GOLD,
             align=PP_ALIGN.CENTER, bold=True)
    add_text(s, l2, card_top + Inches(1.5), card_w, Inches(0.6),
             "RM 280/month · 4 doctors · 6 nurses · analytics",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # Contact block — navy band
    band_top = Inches(5.3)
    add_rect(s, 0, band_top, SLIDE_W, Inches(1.6), NAVY)
    add_text(s, 0, band_top + Inches(0.25), SLIDE_W, Inches(0.5),
             "Hi Kanan team — I'd like to enquire.",
             size=18, bold=True, color=WARM_WHITE, italic=True,
             align=PP_ALIGN.CENTER)
    add_text(s, 0, band_top + Inches(0.85), SLIDE_W, Inches(0.4),
             "WhatsApp +60 12-347 8126   ·   hello@kanan.my   ·   kanan.my",
             size=14, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
    add_speaker_note(s, "Don't ask for the sale here. Ask if they have 30 seconds to try the demo on their phone right now.")


# ─── Build it ─────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        s1_cover,
        s2_pain,
        s3_meet_kanan,
        s4_product_glance,
        s5_patient,
        s6_nurse,
        s7_doctor,
        s8_owner_dash,
        s9_terminal,          # NEW — shared kiosk lockscreen
        s10_room_flow,        # NEW — Premium check-in / check-out flow
        s11_staff_cards,      # NEW — staff cards on /owner/staff
        s12_approvals_audit,  # NEW — Approvals + Audit log
        s13_standard,
        s14_premium,
        s15_comparison,
        s16_how_we_work,
        s17_trust,
        s18_try_it,
    ]
    total = len(builders)
    builders[0](prs)
    for i, fn in enumerate(builders[1:], start=2):
        fn(prs, total)

    out = Path(__file__).parent / "Kanan_Clinic_Booking_Sales_Deck.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")


if __name__ == "__main__":
    build()
