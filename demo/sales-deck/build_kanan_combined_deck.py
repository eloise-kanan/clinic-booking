"""Generate the Kanan clinic-booking combined demo + marketing deck.

ONE deck (per latest user direction): walks through the product with real
mock screenshots, showcases Premium-only features in the middle, top-up
scenario, and Standard-vs-Premium comparison at the back.

Output:  demo/sales-deck/Kanan_Clinic_Booking_Demo_Deck.pptx
Run:     python3 demo/sales-deck/build_kanan_combined_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Kanan brand ─────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_LIGHT  = RGBColor(0x2E, 0x43, 0x74)
GOLD        = RGBColor(0xC9, 0xA2, 0x27)
GOLD_LIGHT  = RGBColor(0xE3, 0xC7, 0x6A)
WARM_WHITE  = RGBColor(0xF4, 0xF1, 0xEA)
GREY        = RGBColor(0x6B, 0x72, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x05, 0x96, 0x69)
RED         = RGBColor(0xC0, 0x39, 0x2B)
AMBER       = RGBColor(0xB4, 0x53, 0x09)

HEAD_FONT = "Helvetica"
BODY_FONT = "Helvetica"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.5)

SCREENSHOTS = Path(__file__).parent / "screenshots"


# ─── Helpers ──────────────────────────────────────────────────────────────

def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, line_color=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill_solid(s, color)
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_w:
            s.line.width = line_w
    return s


def add_text(slide, l, t, w, h, text_s, *, size=14, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text_s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return box


def add_bullets(slide, l, t, w, h, items, *, size=14, color=NAVY,
                bullet_color=GOLD, line_gap_pt=8, bullet_char="●"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(line_gap_pt)
        b = p.add_run()
        b.text = f"{bullet_char}  "
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        b.font.name = BODY_FONT
        rr = p.add_run()
        rr.text = item
        rr.font.size = Pt(size)
        rr.font.color.rgb = color
        rr.font.name = BODY_FONT


def add_note(slide, text_s):
    slide.notes_slide.notes_text_frame.text = text_s


def page_chrome(prs, page, total, title, lead=None, *, tier_pill=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WARM_WHITE)
    add_rect(s, MARGIN, Inches(0.55), Inches(0.4), Inches(0.04), GOLD)
    add_text(s, MARGIN, Inches(0.65), SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    if lead:
        add_text(s, MARGIN, Inches(1.45), SLIDE_W - 2 * MARGIN - Inches(2), Inches(0.5),
                 lead, size=14, italic=True, color=NAVY_LIGHT)
    if tier_pill:
        pill_bg = GOLD if tier_pill == "Premium" else NAVY_LIGHT
        pill_l = SLIDE_W - MARGIN - Inches(1.4)
        add_rect(s, pill_l, Inches(0.75), Inches(1.4), Inches(0.4), pill_bg, line_color=None)
        add_text(s, pill_l, Inches(0.75), Inches(1.4), Inches(0.4),
                 f"{tier_pill} only", size=11, bold=True, color=NAVY if tier_pill == "Premium" else WARM_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Footer
    add_text(s, MARGIN, SLIDE_H - Inches(0.35), Inches(7), Inches(0.3),
             "Kanan Digital Enterprise · kanan.my", size=9, color=GREY)
    add_text(s, SLIDE_W - MARGIN - Inches(1), SLIDE_H - Inches(0.35),
             Inches(1), Inches(0.3), f"{page} / {total}",
             size=9, color=GREY, align=PP_ALIGN.RIGHT)
    return s


def add_screenshot(slide, l, t, w, h, png_name, *, border=True):
    """Add a screenshot from the screenshots dir, sized to fit the box."""
    path = SCREENSHOTS / png_name
    if not path.exists():
        # Fallback to placeholder if mock generation hasn't run
        box = add_rect(slide, l, t, w, h, WARM_WHITE, line_color=GOLD_LIGHT)
        box.line.width = Pt(1.5)
        add_text(slide, l, t, w, h, f"[ {png_name} ]",
                 size=10, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    pic = slide.shapes.add_picture(str(path), l, t, w, h)
    if border:
        pic.line.color.rgb = NAVY_LIGHT
        pic.line.width = Pt(0.75)


# ─── Slides ───────────────────────────────────────────────────────────────

def s_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WARM_WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.18), GOLD)
    add_text(s, 0, Inches(2.0), SLIDE_W, Inches(0.8),
             "kanan", size=80, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, 0, Inches(3.0), SLIDE_W, Inches(0.5),
             "your trusted right hand", size=18, italic=True,
             color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.0), SLIDE_W, Inches(0.6),
             "A booking system built for Malaysian dental clinics",
             size=22, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.7), SLIDE_W, Inches(0.4),
             "by Kanan Digital Enterprise",
             size=13, color=GREY, align=PP_ALIGN.CENTER)
    add_rect(s, 0, SLIDE_H - Inches(0.5), SLIDE_W, Inches(0.5), NAVY)
    add_text(s, 0, SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.4),
             "kanan.my   ·   hello@kanan.my   ·   WhatsApp +60 12-347 8126",
             size=11, color=WARM_WHITE, align=PP_ALIGN.CENTER)


def s_pain(prs, total):
    s = page_chrome(prs, 2, total, "Short-staffed and stretched thin?",
                    "Most days end the same way at a small dental clinic.")
    add_bullets(s, MARGIN, Inches(2.2), SLIDE_W - 2 * MARGIN, Inches(4.5), [
        "50+ WhatsApp messages a day — just to confirm and remind appointments",
        "A no-show wastes a 1-hour slot, and the patient never paid a deposit",
        "The owner has no idea who's busy and who's coasting",
        "Patient records: half in a notebook, half in a Sheet, half in someone's head",
        "When anyone takes leave, the schedule unravels for a week",
    ], size=17, line_gap_pt=12)
    add_note(s, "Let them nod. Don't oversell. Ask: 'Which of these sounds like your Tuesday?'")


def s_meet_kanan(prs, total):
    s = page_chrome(prs, 3, total, "We've been there. So we built the right hand.")
    add_text(s, MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(0.7),
             "Kanan (\"KAH-nahn\") means right in Malay — as in right hand. "
             "The dependable hand on the right side of your business.",
             size=15, italic=True, color=NAVY_LIGHT)
    card_w = Inches(5.0)
    card_h = Inches(2.5)
    card_top = Inches(3.6)
    gap = Inches(0.3)
    left_start = (SLIDE_W - (card_w * 2 + gap)) // 2
    for i, (name, role) in enumerate([
        ("[Founder name 1]", "Co-founder · [role]"),
        ("[Founder name 2]", "Co-founder · [role]"),
    ]):
        l = left_start + i * (card_w + gap)
        add_rect(s, l, card_top, card_w, card_h, WHITE, line_color=NAVY)
        add_text(s, l, card_top + Inches(0.5), card_w, Inches(0.5),
                 "[photo]", size=12, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER)
        add_text(s, l, card_top + Inches(1.3), card_w, Inches(0.5),
                 name, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, l, card_top + Inches(1.85), card_w, Inches(0.4),
                 role, size=12, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(6.4), SLIDE_W - 2 * MARGIN, Inches(0.5),
             "A Malaysian software studio. We build for the businesses around us — not for venture capital.",
             size=13, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


def s_glance(prs, total):
    s = page_chrome(prs, 4, total, "One product. The whole booking lifecycle.")
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) // 3
    cols = [
        ("📱", "Patient self-books", "via your /book page. EN · 中文 · BM."),
        ("💬", "Nurse runs the day", "Pending queue, WhatsApp, reminders, recall."),
        ("📊", "Owner sees everything", "Bookings, performance, audit log."),
    ]
    for i, (icon, h, body) in enumerate(cols):
        l = MARGIN + i * (col_w + Inches(0.3))
        add_rect(s, l, Inches(2.2), col_w, Inches(2.5), WHITE, line_color=NAVY_LIGHT)
        add_text(s, l, Inches(2.35), col_w, Inches(0.7), icon,
                 size=36, align=PP_ALIGN.CENTER)
        add_text(s, l, Inches(3.2), col_w, Inches(0.5),
                 h, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, l, Inches(3.85), col_w, Inches(0.7),
                 body, size=12, color=GREY, align=PP_ALIGN.CENTER)
    # Sub-message
    add_text(s, MARGIN, Inches(5.2), SLIDE_W - 2 * MARGIN, Inches(0.5),
             "All three views, one shared database. No double entry, no copy-paste between systems.",
             size=14, color=NAVY, italic=True, align=PP_ALIGN.CENTER)


def s_book(prs, total):
    s = page_chrome(prs, 5, total, "Patient self-books, no phone tag.",
                    "Identity first, then they pick book / reschedule / cancel.")
    add_bullets(s, MARGIN, Inches(2.0), Inches(5.5), Inches(4.5), [
        "Mobile-first form. 4 taps to a request.",
        "Language toggle: EN · 中文 · BM (name input stays Romanized for IC)",
        "5-week future window. Real-time slot availability — no fake slots.",
        "Submissions land in your nurse's queue, never directly on the calendar",
        "Branded with your clinic's logo + colours, not Kanan's",
    ], size=13, line_gap_pt=8)
    add_screenshot(s, Inches(7), Inches(2.0), Inches(5.8), Inches(5.0), "mock_book.png")


def s_overview(prs, total):
    s = page_chrome(prs, 6, total, "Owner dashboard — a clinic you can actually see.",
                    "This week + today's bookings + recent activity, one screen.")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0),
                   "mock_overview_premium.png")


def s_staff_seats(prs, total):
    s = page_chrome(prs, 7, total, "Staff management with hard seat caps.",
                    "Your plan defines what you can grow to. Top-up out-of-band when you need more.")
    add_bullets(s, MARGIN, Inches(2.0), Inches(5.0), Inches(4.0), [
        "Live seat meters show capacity at a glance",
        "Server-side enforcement — adding a 3rd doctor on Standard fails with a clear message",
        "Deactivated staff don't consume seats — leave handling is free",
        "Top-up via WhatsApp: 'Add 2 more nurses' — no forced tier upgrade",
    ], size=13, line_gap_pt=8)
    add_screenshot(s, Inches(7), Inches(2.0), Inches(5.8), Inches(5.0),
                   "mock_staff_standard.png")


def s_pending(prs, total):
    s = page_chrome(prs, 8, total, "Pending queue — one click to confirm.",
                    "Patient requests land here for nurse review. Never directly on the calendar.")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0), "mock_pending.png")


def s_reminders(prs, total):
    s = page_chrome(prs, 9, total, "Tomorrow's reminders, pre-written.",
                    "Click 'Send' — WhatsApp opens with the message + slot already in it.")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0), "mock_reminders.png")
    add_note(s, "Highlight: nurse audit trail visible — owner can verify who sent what.")


def s_recalls(prs, total):
    s = page_chrome(prs, 10, total, "6-month recall — automatically surfaced.",
                    "Every patient overdue for a checkup appears here, sorted by how overdue.")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0), "mock_recalls.png")
    add_note(s, "Recall is the most reliable revenue lever in dentistry. Until now, it was a nurse's memory + a sticky note.")


def s_templates(prs, total):
    s = page_chrome(prs, 11, total, "WhatsApp templates — owner-editable.",
                    "6 templates: check, confirm, reject, cancel, reminder, recall.")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0), "mock_templates.png")


def s_duty_std(prs, total):
    s = page_chrome(prs, 12, total, "Duty calendar — Standard tier",
                    "Doctor schedules at a glance. Leaves shown in red.")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0),
                   "mock_duty_standard.png")


def s_premium_intro(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(3.5), SLIDE_W, Inches(0.06), GOLD)
    add_text(s, 0, Inches(2.3), SLIDE_W, Inches(0.7),
             "What you unlock with Premium",
             size=44, bold=True, color=WARM_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(3.8), SLIDE_W, Inches(0.5),
             "Nurse duty · doctor + nurse performance · utilization · audit log",
             size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.6), SLIDE_W, Inches(0.4),
             "+ larger seat caps · roadmap items below",
             size=13, color=WARM_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, SLIDE_H - Inches(0.35), Inches(7), Inches(0.3),
             "Kanan Digital Enterprise · kanan.my", size=9, color=GOLD_LIGHT)
    add_text(s, SLIDE_W - MARGIN - Inches(1), SLIDE_H - Inches(0.35),
             Inches(1), Inches(0.3), f"13 / {total}",
             size=9, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)


def s_duty_premium(prs, total):
    s = page_chrome(prs, 14, total, "Duty calendar — Premium",
                    "Both doctors AND nurses on one calendar. Role badges for clarity.",
                    tier_pill="Premium")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0),
                   "mock_duty_premium.png")


def s_doctor_perf(prs, total):
    s = page_chrome(prs, 15, total, "Doctor performance — who's earning your chair-time.",
                    "Per-doctor attendance + no-show metrics, last 7 / 30 / 90 / 180 days.",
                    tier_pill="Premium")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0),
                   "mock_doctor_perf.png")
    add_note(s, "Attendance % colour-coded: green ≥85, amber ≥70, red below. Surfaces underperformance fast.")


def s_nurse_perf(prs, total):
    s = page_chrome(prs, 16, total, "Nurse performance — derived from the audit log.",
                    "Every approval, reminder, recall and attendance mark, attributed.",
                    tier_pill="Premium")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0),
                   "mock_nurse_perf.png")
    add_note(s, "Owners love this. They finally know who's actually working vs who's coasting.")


def s_utilization(prs, total):
    s = page_chrome(prs, 17, total, "Chair utilization — when are you actually full?",
                    "Heatmap of confirmed bookings by weekday × hour.",
                    tier_pill="Premium")
    add_screenshot(s, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(5.0),
                   "mock_utilization.png")
    add_note(s, "Use this to decide working-hour changes or which doctor needs a Saturday shift.")


def s_topup(prs, total):
    s = page_chrome(prs, 18, total, "Seat top-ups — when you grow within your tier",
                    "You don't have to upgrade to add one more doctor or nurse.")
    # 3-card layout
    card_w = Inches(4.0)
    card_h = Inches(4.0)
    card_top = Inches(2.2)
    gap = Inches(0.2)
    left_start = (SLIDE_W - (card_w * 3 + gap * 2)) // 2

    cases = [
        ("Hit your cap", "You try to add a 3rd doctor while on Standard (2-doctor cap).",
         "Server rejects with: \"Seat limit reached — your plan allows 2 doctors.\""),
        ("Tell us on WhatsApp", "One message: \"Hi Kanan, please add 1 more doctor seat.\"",
         "We respond within hours during MY business time. No paperwork."),
        ("Pay only for what you add", "Per-seat top-up: RM 30/mo per extra doctor, RM 20/mo per nurse.",
         "Far cheaper than forcing a full Premium upgrade."),
    ]
    for i, (h, body, extra) in enumerate(cases):
        l = left_start + i * (card_w + gap)
        add_rect(s, l, card_top, card_w, card_h, WHITE, line_color=NAVY_LIGHT)
        # Big number
        add_text(s, l + Inches(0.3), card_top + Inches(0.2), Inches(0.6), Inches(0.6),
                 str(i + 1), size=36, bold=True, color=GOLD, font="Georgia")
        add_text(s, l + Inches(0.3), card_top + Inches(0.95), card_w - Inches(0.6), Inches(0.5),
                 h, size=15, bold=True, color=NAVY)
        add_text(s, l + Inches(0.3), card_top + Inches(1.5), card_w - Inches(0.6), Inches(1.5),
                 body, size=12, color=NAVY)
        add_text(s, l + Inches(0.3), card_top + Inches(2.8), card_w - Inches(0.6), Inches(1.0),
                 extra, size=11, italic=True, color=GREY)
    add_text(s, MARGIN, Inches(6.5), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "💡 Top-up pricing is indicative — confirm during commercial discussion.",
             size=11, italic=True, color=AMBER, align=PP_ALIGN.CENTER)


def s_comparison(prs, total):
    s = page_chrome(prs, 19, total, "Standard vs Premium — at a glance",
                    "Both single-clinic. Franchise is parked until you ask.")
    headers = ["", "Standard", "Premium"]
    rows = [
        ("Monthly price",                          "RM 150",  "RM 280"),
        ("Owner / Doctor / Nurse seats",           "1 / 2 / 3", "1 / 4 / 6"),
        ("Per-seat top-up available",              "✓",       "✓"),
        ("Patient self-book + WhatsApp templates", "✓",       "✓"),
        ("Recall reminders",                       "✓",       "✓"),
        ("Duty calendar — doctors",                "✓",       "✓"),
        ("Duty calendar — nurses",                 "—",       "✓"),
        ("Doctor performance analytics",           "—",       "✓"),
        ("Nurse performance analytics",            "—",       "✓"),
        ("Chair utilization heatmap",              "—",       "✓"),
        ("Audit log",                              "—",       "✓"),
        ("CSV backup + email backup",              "✓",       "✓"),
        ("Branding (logo + colours)",              "✓",       "✓"),
    ]
    table_left = MARGIN + Inches(0.5)
    table_top = Inches(2.0)
    col_w = [Inches(5.5), Inches(3.0), Inches(3.0)]
    row_h = Inches(0.32)
    # Header
    for i, h in enumerate(headers):
        x = table_left + sum(col_w[:i], Inches(0))
        add_rect(s, x, table_top, col_w[i], row_h, NAVY)
        add_text(s, x, table_top, col_w[i], row_h, h,
                 size=12, bold=True, color=WARM_WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows):
        y = table_top + row_h * (ri + 1)
        bg = WARM_WHITE if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            x = table_left + sum(col_w[:ci], Inches(0))
            add_rect(s, x, y, col_w[ci], row_h, bg)
            is_check = val == "✓"
            is_dash = val == "—"
            color = GREEN if is_check else (GREY if is_dash else NAVY)
            add_text(s, x + Inches(0.12), y, col_w[ci] - Inches(0.24), row_h, val,
                     size=11,
                     bold=ci == 0 or is_check,
                     color=color,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
    add_note(s, "Let them stare for 10 seconds. Don't walk through line by line.")


def s_how_we_work(prs, total):
    s = page_chrome(prs, 20, total, "Start to finish — no surprises")
    steps = [
        ("1", "Talk to our team", "WhatsApp or quick call."),
        ("2", "Schedule a discussion", "Longer working session, your pace."),
        ("3", "Workflow & framework", "Map ops + where software fits."),
        ("4", "Proposal & strategic review", "1-page scope, timeline, price."),
        ("5", "Demo", "Working prototype on your real data."),
        ("6", "Commercial discussion", "Full commercial agreed."),
        ("7", "Setup & training", "Go-live, team trained."),
        ("8", "Done — and still here", "Support continues. No ghosting."),
    ]
    grid_top = Inches(1.9)
    cell_w = (SLIDE_W - 2 * MARGIN - Inches(0.3) * 3) // 4
    cell_h = Inches(2.4)
    for i, (n, title, body) in enumerate(steps):
        row, col = divmod(i, 4)
        x = MARGIN + col * (cell_w + Inches(0.3))
        y = grid_top + row * (cell_h + Inches(0.3))
        add_rect(s, x, y, cell_w, cell_h, WHITE, line_color=NAVY_LIGHT)
        add_rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.04), GOLD)
        add_text(s, x + Inches(0.25), y + Inches(0.4), cell_w - Inches(0.5), Inches(0.6),
                 n, size=28, bold=True, color=NAVY, font="Georgia")
        add_text(s, x + Inches(0.25), y + Inches(1.05), cell_w - Inches(0.5), Inches(0.5),
                 title, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), y + Inches(1.55), cell_w - Inches(0.5), Inches(0.8),
                 body, size=10, color=GREY)


def s_trust(prs, total):
    s = page_chrome(prs, 21, total, "Your data stays yours",
                    "Trust before technology.")
    add_bullets(s, MARGIN, Inches(2.3), SLIDE_W - 2 * MARGIN, Inches(4), [
        "Hosted in Malaysia wherever possible (Supabase + Vercel)",
        "PDPA-aware — patient data export any time, schema visible on request",
        "We don't train external AI models on your operational data",
        "CSV backup any day. Daily auto-email backup if you want it.",
        "Sterilization log, X-rays, treatment notes stay in your EMR — not in our system",
    ], size=14, line_gap_pt=10)
    add_text(s, MARGIN, Inches(6.3), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Kanan Digital Enterprise (SSM Reg: 2026060600274 / AS0516554-D) · Kuala Lumpur, Malaysia",
             size=11, italic=True, color=GREY)


def s_try(prs, total):
    s = page_chrome(prs, 22, total, "Try the demos. Set up in 4 weeks.")
    card_w = Inches(5.5)
    card_h = Inches(2.3)
    card_top = Inches(2.0)
    gap = Inches(0.3)
    left_start = (SLIDE_W - (card_w * 2 + gap)) // 2

    # Standard demo
    add_rect(s, left_start, card_top, card_w, card_h, WHITE, line_color=NAVY_LIGHT)
    add_rect(s, left_start, card_top, card_w, Inches(0.15), NAVY_LIGHT)
    add_text(s, left_start, card_top + Inches(0.4), card_w, Inches(0.5),
             "Standard demo", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, left_start, card_top + Inches(1.0), card_w, Inches(0.5),
             "standard-demo.kanan.my", size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, left_start, card_top + Inches(1.5), card_w, Inches(0.6),
             "RM 150/month · 2 doctors · 3 nurses",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # Premium demo
    l2 = left_start + card_w + gap
    add_rect(s, l2, card_top, card_w, card_h, WHITE, line_color=GOLD)
    add_rect(s, l2, card_top, card_w, Inches(0.15), GOLD)
    add_text(s, l2, card_top + Inches(0.4), card_w, Inches(0.5),
             "Premium demo", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, l2, card_top + Inches(1.0), card_w, Inches(0.5),
             "premium-demo.kanan.my", size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, l2, card_top + Inches(1.5), card_w, Inches(0.6),
             "RM 280/month · 4 doctors · 6 nurses · analytics",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # Contact band
    band_top = Inches(5.0)
    add_rect(s, 0, band_top, SLIDE_W, Inches(2.0), NAVY)
    add_text(s, 0, band_top + Inches(0.3), SLIDE_W, Inches(0.5),
             "Hi Kanan team — I'd like to enquire.",
             size=18, bold=True, color=WARM_WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, 0, band_top + Inches(1.0), SLIDE_W, Inches(0.4),
             "WhatsApp +60 12-347 8126   ·   hello@kanan.my   ·   kanan.my",
             size=14, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, 0, band_top + Inches(1.5), SLIDE_W, Inches(0.3),
             "Kanan Digital Enterprise — your trusted right hand",
             size=11, italic=True, color=WARM_WHITE, align=PP_ALIGN.CENTER)
    add_note(s, "Don't close on this slide. Ask if they have 30 seconds to try the demo on their phone right now.")


# ─── Build ───────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        s_cover,
        s_pain,
        s_meet_kanan,
        s_glance,
        s_book,
        s_overview,
        s_staff_seats,
        s_pending,
        s_reminders,
        s_recalls,
        s_templates,
        s_duty_std,
        s_premium_intro,
        s_duty_premium,
        s_doctor_perf,
        s_nurse_perf,
        s_utilization,
        s_topup,
        s_comparison,
        s_how_we_work,
        s_trust,
        s_try,
    ]
    total = len(builders)
    builders[0](prs)
    for i, fn in enumerate(builders[1:], start=2):
        fn(prs, total)

    out = Path(__file__).parent / "Kanan_Clinic_Booking_Demo_Deck.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")


if __name__ == "__main__":
    build()
