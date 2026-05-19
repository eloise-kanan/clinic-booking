"""Generate the Goodcare Dental booking system demo deck.

Creates ./demo/Goodcare_Demo_Walkthrough.pptx with one slide per step in
the booking lifecycle (patient submits → ... → nurse marks attended →
owner reviews audit log).

Each slide has:
- A bold title
- 3-5 main bullet points (what's on the screen)
- A "Speaker note" with what to say out loud during the demo
- A placeholder text box where you can paste a screenshot later

Run:  python3 demo/build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Brand colors
BRAND = RGBColor(0x0D, 0x94, 0x88)        # teal
BRAND_DARK = RGBColor(0x0F, 0x76, 0x6E)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x78, 0x71, 0x6C)
PAPER = RGBColor(0xFA, 0xFA, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xB4, 0x53, 0x09)
RED = RGBColor(0xC0, 0x39, 0x2B)


def add_text_box(slide, left, top, width, height, text, *,
                 font_size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Helvetica"
    return box


def add_bulleted_list(slide, left, top, width, height, bullets,
                      *, font_size=14, color=INK, indent_color=BRAND):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet character
        bullet_run = p.add_run()
        bullet_run.text = "● "
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.color.rgb = indent_color
        bullet_run.font.bold = True
        # actual text
        text_run = p.add_run()
        text_run.text = item
        text_run.font.size = Pt(font_size)
        text_run.font.color.rgb = color
        text_run.font.name = "Helvetica"
        p.space_after = Pt(8)
    return box


def add_brand_bar(slide, slide_width):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()


def add_footer(slide, slide_width, slide_height, page_num, total_pages):
    add_text_box(
        slide,
        Inches(0.5),
        slide_height - Inches(0.4),
        Inches(6),
        Inches(0.3),
        "Goodcare Dental — Booking System Walkthrough",
        font_size=9,
        color=MUTED,
    )
    add_text_box(
        slide,
        slide_width - Inches(1.5),
        slide_height - Inches(0.4),
        Inches(1),
        Inches(0.3),
        f"{page_num} / {total_pages}",
        font_size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_screenshot_placeholder(slide, left, top, width, height, label):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = PAPER
    rect.line.color.rgb = MUTED
    rect.line.width = Pt(0.75)
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {label} ]"
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED
    run.font.italic = True
    run.font.name = "Helvetica"


def add_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_step_pill(slide, left, top, step_num, total_steps):
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.4), Inches(0.32)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = BRAND
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"STEP {step_num} OF {total_steps}"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Helvetica"


# ---- Build the deck ----

prs = Presentation()
prs.slide_width = Inches(13.33)   # 16:9
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# Slide layout 6 = blank
BLANK = prs.slide_layouts[6]

# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_brand_bar(slide, SW)

# Background tint
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.18), SW, SH - Inches(0.18)
)
bg.fill.solid()
bg.fill.fore_color.rgb = PAPER
bg.line.fill.background()

add_text_box(
    slide, Inches(0.8), Inches(2.4), Inches(11), Inches(0.5),
    "Goodcare Dental",
    font_size=20, color=BRAND_DARK, bold=True,
)
add_text_box(
    slide, Inches(0.8), Inches(2.95), Inches(11), Inches(1),
    "Online Booking System",
    font_size=44, color=INK, bold=True,
)
add_text_box(
    slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.6),
    "End-to-end walkthrough — from patient request to completed visit",
    font_size=18, color=MUTED,
)
add_text_box(
    slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
    "Demo · May 2026",
    font_size=12, color=MUTED,
)

add_speaker_notes(slide,
    "Open with a 30-second framing: 'Today I want to walk you through how the system handles "
    "a real patient — from the moment they submit a booking online, to the day they walk into "
    "the clinic. Every step is something your team will use day-to-day.'\n\n"
    "Don't dive into tech. Keep it about what the patient sees and what the staff does."
)

# ============================================================
# Slide 2: The lifecycle (visual flow)
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_brand_bar(slide, SW)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "The patient lifecycle, in 9 steps", font_size=28, bold=True, color=INK)

steps = [
    ("1", "Patient submits", "Online form"),
    ("2", "Nurse sees pending", "Notification"),
    ("3", "Check via WhatsApp", "Verify with patient"),
    ("4", "Patient confirms", "Reply YES"),
    ("5", "Nurse approves", "One click"),
    ("6", "Confirmation sent", "Auto-fills WhatsApp"),
    ("7", "Day-before reminder", "Auto-fills WhatsApp"),
    ("8", "Patient attends", "Mark as attended"),
    ("9", "Audit trail", "Owner reviews"),
]

box_w = Inches(1.25)
box_h = Inches(1.7)
gap = Inches(0.07)
total_w = box_w * 9 + gap * 8
start_x = (SW - total_w) / 2
y = Inches(2.5)
for i, (num, title, sub) in enumerate(steps):
    x = start_x + (box_w + gap) * i
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BRAND
    rect.line.width = Pt(1.5)

    add_text_box(slide, x, y + Inches(0.15), box_w, Inches(0.35),
                 num, font_size=18, color=BRAND, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.55), box_w - Inches(0.1),
                 Inches(0.65), title, font_size=11, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), y + Inches(1.15), box_w - Inches(0.1),
                 Inches(0.5), sub, font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.4),
             "Each step takes seconds. We'll go through them one at a time.",
             font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.4),
             "Nothing is automated where judgement matters. The nurse always controls when a patient is contacted, "
             "approved, or marked as attended.",
             font_size=12, color=INK, align=PP_ALIGN.CENTER)

add_footer(slide, SW, SH, 2, 13)
add_speaker_notes(slide,
    "This is the map of the demo. Nine steps. Most happen in seconds.\n\n"
    "Key point to land: 'The system never sends a message on its own. Every WhatsApp goes through "
    "your nurse's own phone, with the message pre-filled. So patients hear from a real person at "
    "your clinic, not a bot.'"
)

# ============================================================
# Helper for step slides
# ============================================================
TOTAL_STEPS = 9


def step_slide(step_num, title, what_user_sees, key_points, screenshot_label, speaker_text,
               page_num):
    slide = prs.slides.add_slide(BLANK)
    add_brand_bar(slide, SW)

    # Step pill top-left
    add_step_pill(slide, Inches(0.8), Inches(0.55), step_num, TOTAL_STEPS)

    # Title
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
                 title, font_size=28, bold=True, color=INK)
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
                 what_user_sees, font_size=14, color=MUTED)

    # Left column: bullets
    add_text_box(slide, Inches(0.8), Inches(2.55), Inches(5.5), Inches(0.4),
                 "What's on the screen", font_size=12, bold=True, color=BRAND_DARK)
    add_bulleted_list(slide, Inches(0.8), Inches(2.95), Inches(5.5), Inches(3.5),
                      key_points, font_size=13)

    # Right column: screenshot placeholder
    add_text_box(slide, Inches(6.7), Inches(2.55), Inches(6), Inches(0.4),
                 "Live view", font_size=12, bold=True, color=BRAND_DARK)
    add_screenshot_placeholder(slide, Inches(6.7), Inches(2.95),
                               Inches(6), Inches(3.7), screenshot_label)

    add_footer(slide, SW, SH, page_num, 13)
    add_speaker_notes(slide, speaker_text)


# ============================================================
# Slide 3 — Step 1: Patient submits
# ============================================================
step_slide(
    step_num=1,
    title="Patient submits a booking",
    what_user_sees="Public form at /book — patient has the URL from the clinic's website, "
                   "WhatsApp, or QR code.",
    key_points=[
        "Picks Booking / Reschedule / Cancel at the top",
        "Fills name, nationality, IC or passport, WhatsApp",
        "Selects treatment (Scaling, Root canal, Whitening, Wisdom tooth, or Others)",
        "Chooses doctor and date — only available slots are shown",
        "Submits. Sees: 'Request received. Nurse will contact you within 24h.'",
    ],
    screenshot_label="Screenshot of /book — desktop",
    speaker_text=(
        "This is the only thing the patient ever uses. No login, no app to download.\n\n"
        "Walk through the form: 'Notice the slot picker only shows times when this doctor is "
        "actually available — if Dr. Lim is on leave that day, those dates are blocked. The "
        "treatment they pick determines how long the slot is — 30 min for scaling, 90 for root "
        "canal — so the calendar fills up cleanly.'\n\n"
        "Mention: 'Nationality dropdown sets the WhatsApp country code automatically — they only "
        "type the local number.'"
    ),
    page_num=3,
)

# ============================================================
# Slide 4 — Step 2: Nurse sees pending
# ============================================================
step_slide(
    step_num=2,
    title="Nurse sees the pending request",
    what_user_sees="Nurse logs into /nurse — the Pending queue is the first screen, with a "
                   "live count badge.",
    key_points=[
        "Queue card shows: patient name, WhatsApp number, doctor, slot, treatment reason",
        "First-time patient and 'returning patient · 3 prior visits' badges",
        "WhatsApp send buttons pre-filled with the right template",
        "Approve / Reject buttons inline",
        "Auto-expires after 24h if no action — flag in the system",
    ],
    screenshot_label="Screenshot of /nurse — Pending queue",
    speaker_text=(
        "When the patient submits, this row appears in the nurse's queue with a red badge in the "
        "side nav. No email, no SMS — the nurse just sees it the next time she opens the dashboard.\n\n"
        "Highlight the prior-visits badge: 'A returning patient is flagged so the nurse can "
        "personalise the message.'\n\n"
        "Mention the 24h auto-expire: 'If the nurse forgets, the system marks the request expired "
        "after 24 hours so it doesn't sit there forever.'"
    ),
    page_num=4,
)

# ============================================================
# Slide 5 — Step 3: Check via WhatsApp
# ============================================================
step_slide(
    step_num=3,
    title="Nurse verifies with the patient via WhatsApp",
    what_user_sees="Nurse taps 'Check with patient'. Her own WhatsApp opens with the verification "
                   "message pre-filled.",
    key_points=[
        "Tap opens wa.me link — message is already typed",
        "Goes from her clinic-issued WhatsApp account, not the system's",
        "Template is editable per clinic (Settings → WhatsApp templates)",
        "System records the click: who sent, when",
        "No per-message API cost — uses her existing WhatsApp",
    ],
    screenshot_label="Screenshot of WhatsApp template (mobile mockup)",
    speaker_text=(
        "This is the pivotal design choice. We don't pay Meta for the WhatsApp Business API. "
        "Instead, we use the wa.me deep-link standard — the message is pre-filled in the nurse's "
        "WhatsApp, she just hits Send. No setup, no per-message charges.\n\n"
        "Key point: 'Patients see a real human conversation from your clinic's number — not a "
        "bot, not an unfamiliar number. Trust is preserved.'\n\n"
        "Show the templates page: 'You can edit every template — Check, Confirmation, Reschedule, "
        "Cancellation, Reject, Reminder — and the change applies system-wide immediately.'"
    ),
    page_num=5,
)

# ============================================================
# Slide 6 — Step 4: Patient confirms
# ============================================================
step_slide(
    step_num=4,
    title="Patient replies — the conversation is normal WhatsApp",
    what_user_sees="Patient replies 'YES' (or asks for changes) to the nurse's message — entirely "
                   "outside the app.",
    key_points=[
        "No special patient-side software",
        "Conversation history lives in WhatsApp — what nurses already know",
        "If patient asks to change time/doctor, nurse can update the booking in the dashboard",
        "If patient never replies, nurse can resend or eventually reject",
    ],
    screenshot_label="Mockup: WhatsApp chat thread",
    speaker_text=(
        "There's no patient-facing app. The patient lives in WhatsApp like always. The nurse "
        "lives in the booking dashboard. The two systems talk through one wa.me link.\n\n"
        "Anticipate the question: 'What if the patient wants to change the time?' Answer: 'The "
        "nurse can reschedule the booking from the dashboard — one button — and the wa.me confirmation "
        "uses the new time.'"
    ),
    page_num=6,
)

# ============================================================
# Slide 7 — Step 5: Nurse approves
# ============================================================
step_slide(
    step_num=5,
    title="Nurse approves the booking",
    what_user_sees="Nurse clicks Approve in the Pending queue — booking flips to 'Confirmed'.",
    key_points=[
        "Status changes from 'pending' to 'confirmed'",
        "Doctor's calendar shows the booking as a green block",
        "Slot is locked — no other patient can book this time",
        "Audit log records: who approved, when",
        "If it was a reschedule request, the original booking is auto-cancelled",
    ],
    screenshot_label="Screenshot: Approve action + confirmed booking on calendar",
    speaker_text=(
        "One click. The booking is now official.\n\n"
        "Two things happen behind the scenes: the unique-slot index in the database guarantees no "
        "double-booking is possible — even if two patients try at the exact same second, one gets "
        "the slot and the other sees 'try another time'. And every action is timestamped against "
        "who did it, so there's an audit trail you can review later.\n\n"
        "Demo the calendar view briefly: 'See how it shows up immediately on Dr. Lim's "
        "calendar — patient name, treatment, time slot proportional to its duration.'"
    ),
    page_num=7,
)

# ============================================================
# Slide 8 — Step 6: Confirmation sent
# ============================================================
step_slide(
    step_num=6,
    title="Confirmation message sent",
    what_user_sees="Nurse clicks 'Send confirmation' — WhatsApp opens with the confirmation "
                   "template, she sends it.",
    key_points=[
        "WhatsApp opens with confirmation template pre-filled",
        "Message ends with the clinic's name and arrival instructions",
        "Send tracking: '✓ Confirmation sent by [nurse], 2:43 PM' shows on the booking row",
        "If the nurse forgets to send, the booking row keeps showing 'Send confirmation'",
        "Patient now has a written record of their appointment",
    ],
    screenshot_label="Mockup: confirmation WhatsApp message",
    speaker_text=(
        "Confirmation goes via WhatsApp. The patient now has their appointment in writing — "
        "doctor name, date, time. They can refer back to it any time.\n\n"
        "Highlight the tracking pills: 'On the bookings page you'll see green pills — Confirm "
        "sent, Reminder sent — with timestamps and the nurse's name. So if a patient ever calls "
        "saying 'I never got a confirmation', you can see exactly what happened.'"
    ),
    page_num=8,
)

# ============================================================
# Slide 9 — Step 7: Day-before reminder
# ============================================================
step_slide(
    step_num=7,
    title="The day before — reminder",
    what_user_sees="Nurse opens 'Send reminders' tab (or the bookings table). Tomorrow's "
                   "confirmed appointments are listed.",
    key_points=[
        "Each row: patient, time, doctor, treatment",
        "'Send reminder' button → wa.me opens with reminder template",
        "Already sent? Pill shows '✓ Reminder sent by [nurse], 5:02 PM'",
        "Bulk-send by working through the list",
        "Reminder template editable in Settings → WhatsApp templates",
    ],
    screenshot_label="Screenshot: Send reminders page or bookings table reminder buttons",
    speaker_text=(
        "Reminders cut no-shows in half — that's industry data, not ours.\n\n"
        "The nurse can do this in 5 minutes the evening before: open Send Reminders, see "
        "tomorrow's list, tap each one. WhatsApp queues up, she hits send, moves to the next.\n\n"
        "Mention: 'You can re-send if you didn't reach them the first time. Each send is logged.'"
    ),
    page_num=9,
)

# ============================================================
# Slide 10 — Step 8: Patient arrives, marked attended
# ============================================================
step_slide(
    step_num=8,
    title="Appointment day — patient arrives",
    what_user_sees="Nurse goes to All bookings → Today tab. As patients walk in, she taps "
                   "'✓ Attended'. If someone doesn't come, she taps 'No-show'.",
    key_points=[
        "'Today' quick-filter shows just today's confirmed bookings",
        "Search by IC or patient name when they walk in",
        "✓ Attended: visit count goes up, patient flagged as completed",
        "No-show: tracked separately — useful for spotting repeat offenders",
        "Both states are reversible — Undo button in case of mistake",
    ],
    screenshot_label="Screenshot: Today's bookings with Attended/No-show actions",
    speaker_text=(
        "This is what the receptionist actually does during the day. Patient walks in, she "
        "searches their IC, taps Attended, moves on.\n\n"
        "Why this matters: 'You can now see your no-show rate over time. If a particular patient "
        "no-shows 3 times in a row, you decide whether to require deposit for future bookings. "
        "Without this data, you'd never know.'"
    ),
    page_num=10,
)

# ============================================================
# Slide 11 — Step 9: Audit log
# ============================================================
step_slide(
    step_num=9,
    title="Owner reviews the audit log",
    what_user_sees="Owner navigates to Settings → Audit log. Sees every action logged "
                   "with timestamp, actor, and details.",
    key_points=[
        "Every approve / reject / cancel / override / password reset is logged",
        "Filter by action type, actor, entity type",
        "Click 'Show details' for the before/after JSON of any action",
        "Used for: dispute resolution, accountability, compliance evidence",
        "Required for PDPA — clinic can show 'who accessed what' if asked",
    ],
    screenshot_label="Screenshot: /owner/audit",
    speaker_text=(
        "This is the owner's safety net. Every meaningful action — approve, cancel, reset password, "
        "edit a template — is logged with the staff member's name, the time, and a snapshot of "
        "what changed.\n\n"
        "Two real uses: dispute resolution (patient says 'I cancelled, why was I charged?' — you "
        "can show the audit entry), and PDPA compliance (Malaysian Personal Data Protection Act "
        "requires you to demonstrate access controls and traceability).\n\n"
        "Even if a nurse leaves the clinic, her name stays on every action she did — we never "
        "delete history."
    ),
    page_num=11,
)

# ============================================================
# Slide 12 — Beyond the lifecycle
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_brand_bar(slide, SW)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "What else is in the system", font_size=28, bold=True, color=INK)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "Features the lifecycle didn't show", font_size=14, color=MUTED)

extras = [
    ("Doctor calendars",
     "Each doctor sees only their own schedule. Block out time for lunch, MC, conferences."),
    ("Duty calendar",
     "Month / week view of who's on duty. Approved leaves and custom shifts overlay clearly."),
    ("Leave + shift change requests",
     "Nurse / doctor submits, owner approves. Approved leave automatically blocks the doctor's "
     "booking calendar."),
    ("Multi-role access",
     "Owner / nurse / doctor — each sees only what their role needs. Add or deactivate staff "
     "from the dashboard. Reset passwords."),
    ("WhatsApp templates",
     "All 6 message templates editable per clinic, with placeholder reference and live preview."),
    ("Mobile-friendly",
     "Add to phone Home Screen for one-tap access. Calendar, queue, and reminders work on phone."),
]

x = Inches(0.8)
y = Inches(2.0)
for i, (title, desc) in enumerate(extras):
    row = i // 2
    col = i % 2
    cx = x + Inches(6.1) * col
    cy = y + Inches(1.65) * row
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  cx, cy, Inches(5.9), Inches(1.45))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = MUTED
    rect.line.width = Pt(0.75)
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.15),
                 Inches(5.4), Inches(0.4), title,
                 font_size=14, bold=True, color=BRAND_DARK)
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.55),
                 Inches(5.4), Inches(0.85), desc,
                 font_size=11, color=INK)

add_footer(slide, SW, SH, 12, 13)
add_speaker_notes(slide,
    "Quick survey of features outside the patient lifecycle. Don't dwell — call out one or two "
    "the clinic seems most interested in. The doctor calendar and the leave/shift workflow "
    "usually land best with clinic owners because they see the value immediately.\n\n"
    "If they ask 'what else?' — pull up the duty calendar live and show them. It's a strong "
    "visual."
)

# ============================================================
# Slide 13 — Closing / Q&A
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_brand_bar(slide, SW)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.7),
             "Q&A", font_size=44, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.5),
             "What would you want to change for your clinic?",
             font_size=18, color=MUTED, align=PP_ALIGN.CENTER)

# Three boxes
items = [
    ("Cost", "Free for first 3 months, then RM 80/month\n(less than one no-show)."),
    ("Setup", "We migrate your existing patients\nand doctor schedules in 1 day."),
    ("Risk", "Walk away anytime —\nyour data exports as CSV."),
]
y = Inches(4.0)
box_w = Inches(3.8)
gap = Inches(0.3)
total_w = box_w * 3 + gap * 2
start_x = (SW - total_w) / 2
for i, (title, body) in enumerate(items):
    cx = start_x + (box_w + gap) * i
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  cx, y, box_w, Inches(2.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = BRAND
    rect.line.width = Pt(1.5)
    add_text_box(slide, cx, y + Inches(0.3), box_w, Inches(0.5),
                 title, font_size=18, bold=True, color=BRAND_DARK,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + Inches(0.2), y + Inches(0.95),
                 box_w - Inches(0.4), Inches(1.2),
                 body, font_size=13, color=INK, align=PP_ALIGN.CENTER)

add_footer(slide, SW, SH, 13, 13)
add_speaker_notes(slide,
    "Don't read the boxes. Pause. Look at them. Wait for questions.\n\n"
    "Likely questions:\n"
    "- 'Can we change the colors?' Yes, Tailwind theme — tweak in 5 mins.\n"
    "- 'Does it work offline?' No, requires internet (Vercel hosted).\n"
    "- 'What if WhatsApp goes down?' The booking system still works; nurses just can't send "
    "messages until WhatsApp is back. Bookings, calendar, status changes all keep working.\n"
    "- 'Can patients book without WhatsApp?' Currently they need a WhatsApp number for the nurse "
    "to verify. We could add SMS later if needed.\n"
    "- 'Can we export everything?' Yes, CSV of patients + bookings any time."
)

# ============================================================
# Save
# ============================================================
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "Goodcare_Demo_Walkthrough.pptx")
prs.save(out_path)
print(f"Wrote {out_path}")
print(f"Slides: {len(prs.slides)}")
